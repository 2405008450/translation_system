# -*- coding: utf-8 -*-
"""临时：验证标签合并后 id 数量下降、导出仍正确。用完即删。"""
from io import BytesIO
from pptx import Presentation
from app.services.adapters.pptx_adapter import PptxAdapter
from app.services.adapters.pptx_exporter import PptxExporter
from app.services.file_record_service import _split_format_tagged_text

PATH = r"C:\Users\H\Desktop\测试项目\中翻译\测试文件\新建 PPTX 演示文稿 (2).pptx"
raw = open(PATH, "rb").read()
parsed = PptxAdapter().parse_with_options(raw, filename="x.pptx")
for s in parsed.segments:
    if "⟦" in (s.source_layout_text or ""):
        n = len(set(__import__("re").findall(r"⟦\s*/?\s*(\d+)\s*⟧", s.source_layout_text)))
        print(f"tags={n} layout={s.source_layout_text[:90]}")

# 导出往返（恒等译文，标签保留）
export_segments = []
for s in parsed.segments:
    clean, layout = _split_format_tagged_text(s.source_layout_text or s.source_text)
    export_segments.append({"sentence_id": s.segment_id, "target_text": clean, "target_layout_text": layout})
out = PptxExporter().export(raw, export_segments)
prs = Presentation(BytesIO(out))
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            runs = para.runs
            if len(runs) >= 2 and "可根据工况选用" in "".join(r.text for r in runs):
                print("\n可根据工况选用段落 run 数:", len(runs))
                for r in runs[:8]:
                    print(f"  {r.text[:16]!r} bold={r.font.bold}")
                raise SystemExit
