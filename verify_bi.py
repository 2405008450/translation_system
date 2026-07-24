# -*- coding: utf-8 -*-
"""临时：双语导出保留原文 run 格式 + 追加译文。用完即删。"""
from io import BytesIO
from pptx import Presentation
from app.services.adapters.pptx_adapter import PptxAdapter
from app.services.adapters.pptx_exporter import PptxExporter

PATH = r"C:\Users\H\Desktop\测试项目\中翻译\测试文件\新建 PPTX 演示文稿 (2).pptx"
raw = open(PATH, "rb").read()
parsed = PptxAdapter().parse_with_options(raw, filename="x.pptx")

# 用带标签的版式原文当作“译文”（标签保留），模拟译文
translated = [{"sentence_id": s.segment_id, "target_text": (s.source_layout_text or s.source_text)} for s in parsed.segments]

out = PptxExporter().export(raw, translated, bilingual=True)
prs = Presentation(BytesIO(out))
shown = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            runs = para.runs
            joined = "".join(r.text for r in runs)
            if "聚合物抗裂砂浆是通过" in joined:
                print("=== 双语混排段落 runs ===")
                for r in runs:
                    rpr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    print(f"  {r.text[:22]!r} bold={r.font.bold} u={r.font.underline} sz={r.font.size}")
                # 检查是否有 br
                xml = para._p.xml
                print("has <a:br/>:", "<a:br" in xml or ':br' in xml)
                shown += 1
    if shown:
        break
