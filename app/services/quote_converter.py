"""引号转换工具的核心算法。

从独立的 GUI 工具 (E:\\it\\引号转换\\text_tool_gui.py) 移植而来，
去除了 tkinter GUI，只保留纯字节级别的转换能力，供 HTTP 接口调用。

支持的输入类型：
- 纯文本类：.txt .md .markdown .srt .rtf（按探测到的编码读写）
- 富文本类：.docx .xlsx .pptx（保留原有格式，DOCX 会调整弯引号字体 hint）
- HTML 类：.html .htm（只处理文本节点，跳过 script/style/注释）
"""

from __future__ import annotations

import io
import os
from copy import deepcopy
from typing import Optional, Set, Tuple

from bs4 import BeautifulSoup, Comment, NavigableString
from docx import Document
from docx.oxml.ns import qn
from lxml import etree
from openpyxl import load_workbook
from pptx import Presentation


# =========================
# 支持的格式
# =========================
PLAIN_TEXT_EXTS = {".txt", ".md", ".markdown", ".srt", ".rtf"}
RICH_TEXT_EXTS = {".docx", ".xlsx", ".pptx"}
HTML_EXTS = {".html", ".htm"}
SUPPORTED_EXTS = PLAIN_TEXT_EXTS | RICH_TEXT_EXTS | HTML_EXTS


# =========================
# 引号字符定义
# =========================
HALF_STRAIGHT_DOUBLE = '"'
HALF_STRAIGHT_SINGLE = "'"
FULL_STRAIGHT_DOUBLE = "\uff02"
FULL_STRAIGHT_SINGLE = "\uff07"
CURLY_DOUBLE_LEFT = "\u201c"
CURLY_DOUBLE_RIGHT = "\u201d"
CURLY_SINGLE_LEFT = "\u2018"
CURLY_SINGLE_RIGHT = "\u2019"

HALF_STRAIGHT = {HALF_STRAIGHT_DOUBLE, HALF_STRAIGHT_SINGLE}
FULL_STRAIGHT = {FULL_STRAIGHT_DOUBLE, FULL_STRAIGHT_SINGLE}
CURLY = {CURLY_DOUBLE_LEFT, CURLY_DOUBLE_RIGHT, CURLY_SINGLE_LEFT, CURLY_SINGLE_RIGHT}

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

SCOPE_WIDTHS = {"全部", "半角", "全角"}
SCOPE_SHAPES = {"全部", "直引号", "弯引号"}
TARGET_WIDTHS = {"半角", "全角"}
TARGET_SHAPES = {"弯引号", "直引号"}

# 常见 MIME 映射，用于返回时设置 Content-Type
MIME_TYPES = {
    ".txt": "text/plain; charset=utf-8",
    ".md": "text/markdown; charset=utf-8",
    ".markdown": "text/markdown; charset=utf-8",
    ".srt": "application/x-subrip; charset=utf-8",
    ".rtf": "application/rtf",
    ".html": "text/html; charset=utf-8",
    ".htm": "text/html; charset=utf-8",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

ENGLISH_FONTS = {
    "times new roman", "arial", "calibri", "cambria", "georgia",
    "verdana", "tahoma", "helvetica", "courier new", "consolas",
    "segoe ui", "trebuchet ms", "palatino linotype", "book antiqua",
    "garamond", "century gothic", "franklin gothic", "impact",
}


class QuoteConverterError(ValueError):
    """引号转换过程中的用户可见错误。"""


def validate_options(scope_width: str, scope_shape: str,
                     target_width: str, target_shape: str) -> None:
    if scope_width not in SCOPE_WIDTHS:
        raise QuoteConverterError(f"无效的处理范围宽度：{scope_width}")
    if scope_shape not in SCOPE_SHAPES:
        raise QuoteConverterError(f"无效的处理范围形状：{scope_shape}")
    if target_width not in TARGET_WIDTHS:
        raise QuoteConverterError(f"无效的目标宽度：{target_width}")
    if target_shape not in TARGET_SHAPES:
        raise QuoteConverterError(f"无效的目标形状：{target_shape}")


def get_extension(filename: str) -> str:
    return os.path.splitext(filename)[1].lower()


def is_supported(filename: str) -> bool:
    return get_extension(filename) in SUPPORTED_EXTS


def media_type_for(filename: str) -> str:
    return MIME_TYPES.get(get_extension(filename), "application/octet-stream")


def build_output_filename(original_name: str, scope_width: str, scope_shape: str,
                          target_width: str, target_shape: str, timestamp: str) -> str:
    """按原 GUI 的规则组合输出文件名：{原名}_{源范围}→{目标}_{时间戳}{扩展名}。"""
    if "." in original_name:
        stem, ext = original_name.rsplit(".", 1)
        ext = f".{ext}"
    else:
        stem, ext = original_name, ""
    action = f"{scope_width}{scope_shape}→{target_width}{target_shape}"
    return f"{stem}_{action}_{timestamp}{ext}"


# =========================
# 编码探测
# =========================

def _detect_encoding(raw: bytes) -> str:
    if raw.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    if raw.startswith(b"\xff\xfe") or raw.startswith(b"\xfe\xff"):
        return "utf-16"
    try:
        from charset_normalizer import from_bytes

        best = from_bytes(raw).best()
        if best and best.encoding:
            return best.encoding
    except Exception:
        pass
    return "utf-8"


def _decode(raw: bytes, encoding: str) -> str:
    try:
        return raw.decode(encoding)
    except UnicodeDecodeError:
        return raw.decode(encoding, errors="replace")


# =========================
# 字体判断（DOCX 用）
# =========================

def _is_english_font(font_name: Optional[str]) -> bool:
    if not font_name:
        return False
    return font_name.lower() in ENGLISH_FONTS


def _get_theme_font(doc, theme_ref: str) -> Optional[str]:
    try:
        theme_part = doc.part.package.parts["/word/theme/theme1.xml"]
        theme_xml = etree.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        if "minor" in theme_ref.lower():
            fonts = theme_xml.xpath("//a:fontScheme/a:minorFont", namespaces=ns)
        else:
            fonts = theme_xml.xpath("//a:fontScheme/a:majorFont", namespaces=ns)
        if fonts:
            font_elem = fonts[0]
            if "eastAsia" in theme_ref:
                ea = font_elem.find("a:ea", namespaces=ns)
                if ea is not None:
                    return ea.get("typeface")
            else:
                latin = font_elem.find("a:latin", namespaces=ns)
                if latin is not None:
                    return latin.get("typeface")
    except Exception:
        return None
    return None


def _get_run_font(run, doc=None) -> Optional[str]:
    if run.font.name:
        return run.font.name
    try:
        rPr = run._element.rPr
        if rPr is not None:
            rFonts = rPr.rFonts
            if rFonts is not None:
                for attr in ("w:ascii", "w:hAnsi", "w:eastAsia", "w:cs"):
                    font = rFonts.get(qn(attr))
                    if font:
                        return font
                for attr in ("w:asciiTheme", "w:hAnsiTheme", "w:eastAsiaTheme"):
                    theme_ref = rFonts.get(qn(attr))
                    if theme_ref and doc is not None:
                        theme_font = _get_theme_font(doc, theme_ref)
                        if theme_font:
                            return theme_font
    except Exception:
        pass
    try:
        para = run._element.getparent()
        if para is not None:
            pPr = para.find(qn("w:pPr"))
            if pPr is not None:
                pStyle = pPr.find(qn("w:pStyle"))
                if pStyle is not None:
                    style_id = pStyle.get(qn("w:val"))
                    if style_id and doc is not None:
                        style = doc.styles.get_by_id(style_id, "paragraph")
                        if style and style.font.name:
                            return style.font.name
    except Exception:
        pass
    try:
        if doc is not None:
            default_style = doc.styles["Normal"]
            if default_style.font.name:
                return default_style.font.name
    except Exception:
        pass
    return None


def _get_run_hint(run) -> Optional[str]:
    try:
        rPr = run._element.rPr
        if rPr is not None:
            rFonts = rPr.rFonts
            if rFonts is not None:
                return rFonts.get(qn("w:hint"))
    except Exception:
        return None
    return None


def _is_curly_half_width(run, doc=None) -> bool:
    hint = _get_run_hint(run)
    if hint:
        return hint == "default"
    return _is_english_font(_get_run_font(run, doc))


# =========================
# 核心转换算法
# =========================

def _scope_set(scope_width: str, scope_shape: str, is_curly_half: Optional[bool]) -> Set[str]:
    result: Set[str] = set()
    if scope_shape in ("全部", "直引号"):
        if scope_width in ("全部", "半角"):
            result |= HALF_STRAIGHT
        if scope_width in ("全部", "全角"):
            result |= FULL_STRAIGHT
    if scope_shape in ("全部", "弯引号"):
        if scope_width == "全部":
            result |= CURLY
        elif scope_width == "半角" and is_curly_half:
            result |= CURLY
        elif scope_width == "全角" and not is_curly_half:
            result |= CURLY
    return result


def _convert_text(text: str, scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str,
                  is_curly_half: Optional[bool],
                  open_double: bool = True,
                  open_single: bool = True) -> Tuple[str, Set[int], bool, bool]:
    """执行字符级别的转换。返回 (新文本, 被转换位置集合, 双引号开合状态, 单引号开合状态)。"""

    scope = _scope_set(scope_width, scope_shape, is_curly_half)
    out = []
    converted: Set[int] = set()

    for ch in text:
        if ch not in scope:
            out.append(ch)
            if ch == CURLY_DOUBLE_LEFT:
                open_double = False
            elif ch == CURLY_DOUBLE_RIGHT:
                open_double = True
            elif ch == CURLY_SINGLE_LEFT:
                open_single = False
            elif ch == CURLY_SINGLE_RIGHT:
                open_single = True
            continue

        is_double = ch in (HALF_STRAIGHT_DOUBLE, FULL_STRAIGHT_DOUBLE,
                           CURLY_DOUBLE_LEFT, CURLY_DOUBLE_RIGHT)
        pos = len(out)

        if target_shape == "弯引号":
            if is_double:
                if ch == CURLY_DOUBLE_LEFT:
                    new_ch = CURLY_DOUBLE_LEFT
                    open_double = False
                elif ch == CURLY_DOUBLE_RIGHT:
                    new_ch = CURLY_DOUBLE_RIGHT
                    open_double = True
                else:
                    new_ch = CURLY_DOUBLE_LEFT if open_double else CURLY_DOUBLE_RIGHT
                    open_double = not open_double
            else:
                if ch == CURLY_SINGLE_LEFT:
                    new_ch = CURLY_SINGLE_LEFT
                    open_single = False
                elif ch == CURLY_SINGLE_RIGHT:
                    new_ch = CURLY_SINGLE_RIGHT
                    open_single = True
                else:
                    new_ch = CURLY_SINGLE_LEFT if open_single else CURLY_SINGLE_RIGHT
                    open_single = not open_single
            out.append(new_ch)
            converted.add(pos)
        else:
            if target_width == "半角":
                new_ch = HALF_STRAIGHT_DOUBLE if is_double else HALF_STRAIGHT_SINGLE
            else:
                new_ch = FULL_STRAIGHT_DOUBLE if is_double else FULL_STRAIGHT_SINGLE
            out.append(new_ch)
            if new_ch != ch:
                converted.add(pos)

    return "".join(out), converted, open_double, open_single


# =========================
# 各格式的处理器
# =========================

def _convert_plain_text(raw: bytes, scope_width: str, scope_shape: str,
                        target_width: str, target_shape: str) -> bytes:
    encoding = _detect_encoding(raw)
    text = _decode(raw, encoding)
    new_text, _, _, _ = _convert_text(
        text, scope_width, scope_shape, target_width, target_shape,
        is_curly_half=None,
    )
    try:
        return new_text.encode(encoding)
    except (LookupError, UnicodeEncodeError):
        return new_text.encode("utf-8")


def _set_curly_font_hint(rPr_elem, target_width: str) -> None:
    rFonts = rPr_elem.find(f"{{{W_NS}}}rFonts")
    if rFonts is None:
        rFonts = etree.SubElement(rPr_elem, f"{{{W_NS}}}rFonts")
    if target_width == "半角":
        rFonts.set(qn("w:hint"), "default")
        rFonts.set(qn("w:ascii"), "Times New Roman")
        rFonts.set(qn("w:hAnsi"), "Times New Roman")
    else:
        rFonts.set(qn("w:hint"), "eastAsia")


def _convert_docx(raw: bytes, scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str) -> bytes:
    try:
        doc = Document(io.BytesIO(raw))
    except Exception as exc:  # noqa: BLE001
        raise QuoteConverterError("DOCX 文件无法解析，请确认文件未损坏") from exc

    need_font_adjust = target_shape == "弯引号"

    for para in doc.paragraphs:
        open_double = True
        open_single = True

        if not need_font_adjust:
            for run in para.runs:
                is_curly_half = _is_curly_half_width(run, doc)
                new_text, _, open_double, open_single = _convert_text(
                    run.text, scope_width, scope_shape,
                    target_width, target_shape, is_curly_half,
                    open_double, open_single,
                )
                run.text = new_text
            continue

        p_element = para._element
        original_runs = list(para.runs)
        new_elements = []

        for run in original_runs:
            old_text = run.text
            is_curly_half = _is_curly_half_width(run, doc)
            new_text, converted_pos, open_double, open_single = _convert_text(
                old_text, scope_width, scope_shape,
                target_width, target_shape, is_curly_half,
                open_double, open_single,
            )

            if not converted_pos:
                new_r = deepcopy(run._element)
                for t in new_r.findall(f".//{{{W_NS}}}t"):
                    t.getparent().remove(t)
                t_elem = etree.SubElement(new_r, f"{{{W_NS}}}t")
                t_elem.text = new_text
                if new_text and (new_text[0] == " " or new_text[-1] == " "):
                    t_elem.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
                new_elements.append(new_r)
                continue

            segments = []
            current_chars = []
            current_is_converted: Optional[bool] = None
            for i, ch in enumerate(new_text):
                is_conv = (i in converted_pos) and (ch in CURLY)
                if current_is_converted is not None and is_conv != current_is_converted:
                    segments.append(("".join(current_chars), current_is_converted))
                    current_chars = []
                current_chars.append(ch)
                current_is_converted = is_conv
            if current_chars:
                segments.append((
                    "".join(current_chars),
                    current_is_converted if current_is_converted is not None else False,
                ))

            for seg_text, is_conv in segments:
                new_r = deepcopy(run._element)
                for t in new_r.findall(f".//{{{W_NS}}}t"):
                    t.getparent().remove(t)
                t_elem = etree.SubElement(new_r, f"{{{W_NS}}}t")
                t_elem.text = seg_text
                if seg_text and (seg_text[0] == " " or seg_text[-1] == " "):
                    t_elem.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
                if is_conv:
                    rPr = new_r.find(f"{{{W_NS}}}rPr")
                    if rPr is None:
                        rPr = etree.Element(f"{{{W_NS}}}rPr")
                        new_r.insert(0, rPr)
                    _set_curly_font_hint(rPr, target_width)
                new_elements.append(new_r)

        for run in original_runs:
            p_element.remove(run._element)
        for elem in new_elements:
            p_element.append(elem)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _convert_xlsx(raw: bytes, scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str) -> bytes:
    try:
        wb = load_workbook(io.BytesIO(raw))
    except Exception as exc:  # noqa: BLE001
        raise QuoteConverterError("XLSX 文件无法解析，请确认文件未损坏") from exc

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                value = cell.value
                if isinstance(value, str) and value:
                    new_text, _, _, _ = _convert_text(
                        value, scope_width, scope_shape,
                        target_width, target_shape, is_curly_half=None,
                    )
                    if new_text != value:
                        cell.value = new_text
    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


def _convert_pptx_text_frame(text_frame, scope_width: str, scope_shape: str,
                             target_width: str, target_shape: str) -> None:
    for para in text_frame.paragraphs:
        open_double = True
        open_single = True
        for run in para.runs:
            if not run.text:
                continue
            new_text, _, open_double, open_single = _convert_text(
                run.text, scope_width, scope_shape,
                target_width, target_shape,
                is_curly_half=None,
                open_double=open_double,
                open_single=open_single,
            )
            if new_text != run.text:
                run.text = new_text


def _convert_pptx_shape(shape, scope_width: str, scope_shape: str,
                        target_width: str, target_shape: str) -> None:
    # 组合形状：递归
    if getattr(shape, "shape_type", None) == 6:
        for sub in shape.shapes:
            _convert_pptx_shape(sub, scope_width, scope_shape, target_width, target_shape)
        return
    # 表格
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                _convert_pptx_text_frame(cell.text_frame, scope_width, scope_shape,
                                         target_width, target_shape)
        return
    if getattr(shape, "has_text_frame", False):
        _convert_pptx_text_frame(shape.text_frame, scope_width, scope_shape,
                                 target_width, target_shape)


def _convert_pptx(raw: bytes, scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str) -> bytes:
    try:
        prs = Presentation(io.BytesIO(raw))
    except Exception as exc:  # noqa: BLE001
        raise QuoteConverterError("PPTX 文件无法解析，请确认文件未损坏") from exc

    for slide in prs.slides:
        for shape in slide.shapes:
            _convert_pptx_shape(shape, scope_width, scope_shape, target_width, target_shape)
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _html_font_for_width(target_width: str) -> str:
    if target_width == "半角":
        return "'Times New Roman', Times, serif"
    return "'SimSun', '宋体', serif"


def _html_make_node(soup, text: str, is_conv_curly: bool, font_family: str):
    if is_conv_curly:
        span = soup.new_tag("span", style=f"font-family:{font_family}")
        span.string = text
        return span
    return NavigableString(text)


def _convert_html(raw: bytes, scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str) -> bytes:
    encoding = _detect_encoding(raw)
    html_text = _decode(raw, encoding)

    soup = BeautifulSoup(html_text, "html.parser")
    skip_tags = {"script", "style"}
    wrap_curly = target_shape == "弯引号"
    font_family = _html_font_for_width(target_width)

    for node in list(soup.find_all(string=True)):
        if node.parent and node.parent.name in skip_tags:
            continue
        if isinstance(node, Comment):
            continue
        text = str(node)
        if not text:
            continue
        new_text, converted_pos, _, _ = _convert_text(
            text, scope_width, scope_shape,
            target_width, target_shape, is_curly_half=None,
        )
        if new_text == text and not (wrap_curly and converted_pos):
            continue
        if not wrap_curly or not converted_pos:
            node.replace_with(NavigableString(new_text))
            continue

        parent = node.parent
        new_nodes = []
        buf = []
        buf_is_curly: Optional[bool] = None
        for i, ch in enumerate(new_text):
            is_conv_curly = (i in converted_pos) and (ch in CURLY)
            if buf_is_curly is not None and is_conv_curly != buf_is_curly:
                new_nodes.append(_html_make_node(soup, "".join(buf), buf_is_curly, font_family))
                buf = []
            buf.append(ch)
            buf_is_curly = is_conv_curly
        if buf:
            new_nodes.append(_html_make_node(soup, "".join(buf), bool(buf_is_curly), font_family))

        idx = list(parent.contents).index(node)
        node.extract()
        for n in reversed(new_nodes):
            parent.insert(idx, n)

    try:
        return str(soup).encode(encoding)
    except (LookupError, UnicodeEncodeError):
        return str(soup).encode("utf-8")


# =========================
# 对外统一入口
# =========================

def convert_bytes(data: bytes, filename: str,
                  scope_width: str, scope_shape: str,
                  target_width: str, target_shape: str) -> Tuple[bytes, str]:
    """按 filename 扩展名分派到对应的处理器。

    返回 (转换后字节, MIME Content-Type)。
    """
    validate_options(scope_width, scope_shape, target_width, target_shape)

    ext = get_extension(filename)
    if ext not in SUPPORTED_EXTS:
        supported = ", ".join(sorted(SUPPORTED_EXTS))
        raise QuoteConverterError(f"不支持的文件类型：{ext or '(无扩展名)'}；支持的格式：{supported}")

    if ext in PLAIN_TEXT_EXTS:
        result = _convert_plain_text(data, scope_width, scope_shape, target_width, target_shape)
    elif ext == ".docx":
        result = _convert_docx(data, scope_width, scope_shape, target_width, target_shape)
    elif ext == ".xlsx":
        result = _convert_xlsx(data, scope_width, scope_shape, target_width, target_shape)
    elif ext == ".pptx":
        result = _convert_pptx(data, scope_width, scope_shape, target_width, target_shape)
    else:  # HTML
        result = _convert_html(data, scope_width, scope_shape, target_width, target_shape)

    return result, media_type_for(filename)
