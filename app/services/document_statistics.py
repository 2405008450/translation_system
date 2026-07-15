from __future__ import annotations

import json
import posixpath
import re
import unicodedata
from collections.abc import Iterable
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from threading import Lock
from typing import Any
from zipfile import BadZipFile, ZipFile
from xml.etree import ElementTree as ET

from app.config import get_settings
from app.services.document_match_analysis import normalize_document_match_analysis
from app.services.libreoffice_service import (
    LibreOfficeError,
    compute_libreoffice_document_statistics,
    convert_word_to_docx,
)


APP_PROPERTIES_NS = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
RELATIONSHIPS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
WP_NS = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
VML_NS = "urn:schemas-microsoft-com:vml"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
DGM_GRAPHIC_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"

STATISTIC_FIELDS = {
    "Pages": "pages",
    "Words": "words",
    "Characters": "characters",
    "CharactersWithSpaces": "characters_with_spaces",
    "Paragraphs": "paragraphs",
    "Lines": "lines",
}
STATISTIC_NUMBER_KEYS = (
    "pages",
    "words",
    "non_asian_words",
    "asian_characters",
    "characters",
    "characters_with_spaces",
    "paragraphs",
    "lines",
    "internal_repeated_words",
    "internal_repeated_characters",
    "cross_file_repeated_words",
    "cross_file_repeated_characters",
    "image_count",
    "unique_image_count",
    "inline_image_count",
    "floating_image_count",
    "linked_image_count",
    "chart_count",
    "smartart_count",
)
NON_ASIAN_WORD_PATTERN = re.compile(
    r"[A-Za-z0-9]+(?:[.,:/_-][A-Za-z0-9]+)*%?|[^\W_\d\s]+(?:[-'][^\W_\d\s]+)*",
    re.UNICODE,
)

_LICENSE_LOCK = Lock()
_LICENSE_STATUS: str | None = None


def compute_docx_statistics(raw_bytes: bytes) -> dict[str, Any]:
    """Compute DOCX statistics with the default OpenXML word-like profile."""
    return compute_word_document_statistics(raw_bytes, "source.docx")


def compute_document_statistics(raw_bytes: bytes, filename: str) -> dict[str, Any]:
    """根据文件类型计算文档统计，同时保持既有 Word 统计口径不变。"""
    suffix = Path(filename or "").suffix.lower()
    if suffix == ".pptx":
        return compute_pptx_document_statistics(raw_bytes)
    return compute_word_document_statistics(raw_bytes, filename)


def compute_word_document_statistics(raw_bytes: bytes, filename: str) -> dict[str, Any]:
    """Compute Word statistics without starting LibreOffice for DOCX files."""
    suffix = Path(filename or "").suffix.lower()
    if suffix == ".docx":
        return _compute_docx_statistics_without_libreoffice(raw_bytes)

    if suffix == ".doc":
        converted_docx = _convert_doc_to_docx(raw_bytes, filename)
        if converted_docx is not None:
            return _compute_docx_statistics_without_libreoffice(converted_docx)

        libreoffice_statistics = _compute_with_libreoffice(raw_bytes, filename)
        if libreoffice_statistics is not None:
            return libreoffice_statistics

    return _build_statistics_payload(
        source="unavailable",
        engine=None,
        include_textboxes_footnotes_endnotes=False,
        license_status=None,
        statistics_profile="unavailable",
        content_scope="unavailable",
    )


def compute_pptx_document_statistics(raw_bytes: bytes) -> dict[str, Any]:
    """按 PPTX 形状文本 + Word 近似口径统计演示文稿。"""
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(raw_bytes))
    except Exception:
        return _build_statistics_payload(
            source="unavailable",
            engine=None,
            include_textboxes_footnotes_endnotes=False,
            license_status=None,
            statistics_profile="unavailable",
            content_scope="unavailable",
        )

    main_texts: list[str] = []
    main_paragraphs = 0
    main_lines = 0
    notes_paragraphs = 0
    notes_lines = 0
    chart_count = 0

    for slide in presentation.slides:
        slide_texts: list[str] = []
        for shape in _iter_ppt_shapes(slide.shapes):
            shape_texts, paragraph_count, line_count, shape_chart_count = _extract_ppt_shape_texts(shape)
            slide_texts.extend(shape_texts)
            main_paragraphs += paragraph_count
            main_lines += line_count
            chart_count += shape_chart_count
        if slide_texts:
            main_texts.append("\n".join(slide_texts))

        notes_text, paragraph_count, line_count = _extract_ppt_notes_text(slide)
        if notes_text:
            notes_paragraphs += paragraph_count
            notes_lines += line_count

    metrics = _count_tool_word_like_texts(main_texts)
    image_count, linked_image_count = _count_pptx_image_references(raw_bytes)
    payload = _build_statistics_payload(
        source="pptx_word_like",
        engine="python-pptx-word-like",
        include_textboxes_footnotes_endnotes=False,
        license_status=None,
        statistics_profile="pptx_shape_word_approx",
        content_scope="slides_shapes_tables_charts",
        statistics_warnings=(
            ["pptx_notes_excluded_from_word_count"]
            if notes_paragraphs or notes_lines
            else []
        ),
    )
    payload.update(
        {
            "pages": len(presentation.slides),
            "words": metrics["words"],
            "non_asian_words": metrics["non_asian_words"],
            "asian_characters": metrics["asian_characters"],
            "characters": metrics["characters"],
            "characters_with_spaces": metrics["characters_with_spaces"],
            "paragraphs": main_paragraphs + notes_paragraphs,
            "lines": main_lines + notes_lines,
            "image_count": image_count,
            "linked_image_count": linked_image_count,
            "chart_count": chart_count,
        }
    )
    return payload


def _convert_doc_to_docx(raw_bytes: bytes, filename: str) -> bytes | None:
    try:
        return convert_word_to_docx(raw_bytes, filename)
    except LibreOfficeError:
        return None


def _compute_with_libreoffice(raw_bytes: bytes, filename: str) -> dict[str, Any] | None:
    try:
        statistics = compute_libreoffice_document_statistics(raw_bytes, filename)
    except LibreOfficeError:
        return None
    statistics.setdefault("statistics_profile", "libreoffice")
    statistics.setdefault("content_scope", "libreoffice_default")
    statistics.setdefault("statistics_warnings", [])
    return statistics


def _compute_docx_statistics_without_libreoffice(raw_bytes: bytes) -> dict[str, Any]:
    openxml_statistics = _compute_with_openxml(raw_bytes)
    if openxml_statistics is not None:
        return openxml_statistics

    aspose_statistics = _compute_with_aspose(raw_bytes)
    if aspose_statistics is not None:
        return aspose_statistics

    cached_statistics = _read_cached_docprops_statistics(raw_bytes)
    if cached_statistics is not None:
        return cached_statistics

    return _build_statistics_payload(
        source="unavailable",
        engine=None,
        include_textboxes_footnotes_endnotes=False,
        license_status=None,
        statistics_profile="unavailable",
        content_scope="unavailable",
    )


def normalize_document_statistics(value: Any) -> dict[str, Any]:
    if isinstance(value, dict):
        if not value:
            return {}
        return _coerce_statistics_payload(value)
    if isinstance(value, str) and value.strip():
        try:
            decoded = json.loads(value)
        except (TypeError, ValueError):
            return {}
        if isinstance(decoded, dict):
            if not decoded:
                return {}
            return _coerce_statistics_payload(decoded)
    return {}


def serialize_document_statistics(value: Any) -> str:
    statistics = normalize_document_statistics(value)
    return json.dumps(statistics, ensure_ascii=False, sort_keys=True)


def _compute_with_aspose(raw_bytes: bytes) -> dict[str, Any] | None:
    try:
        import aspose.words as aw
    except Exception:
        return None

    license_status = _apply_aspose_license(aw)
    try:
        with TemporaryDirectory(prefix="docx-statistics-") as temp_dir:
            docx_path = Path(temp_dir) / "source.docx"
            docx_path.write_bytes(raw_bytes)
            document = aw.Document(str(docx_path))
            document.include_textboxes_footnotes_endnotes_in_stat = False
            document.update_word_count(True)
            properties = document.built_in_document_properties

            payload = _build_statistics_payload(
                source="aspose",
                engine="aspose-words",
                include_textboxes_footnotes_endnotes=False,
                license_status=license_status,
                statistics_profile="aspose_word",
                content_scope="aspose_document",
            )
            payload.update(
                {
                    "pages": _first_int(
                        getattr(document, "page_count", None),
                        getattr(properties, "pages", None),
                    ),
                    "words": _to_optional_int(getattr(properties, "words", None)),
                    "characters": _to_optional_int(getattr(properties, "characters", None)),
                    "characters_with_spaces": _to_optional_int(
                        getattr(properties, "characters_with_spaces", None)
                    ),
                    "paragraphs": _to_optional_int(getattr(properties, "paragraphs", None)),
                    "lines": _to_optional_int(getattr(properties, "lines", None)),
                }
            )
            return payload
    except Exception:
        return None


def _compute_with_openxml(raw_bytes: bytes) -> dict[str, Any] | None:
    try:
        with ZipFile(BytesIO(raw_bytes)) as archive:
            document_xml = archive.read("word/document.xml")
            cached_statistics = _read_cached_docprops_statistics_from_archive(archive)
            relationships = _read_part_relationships(archive, "word/document.xml")
    except (BadZipFile, KeyError):
        return None

    try:
        root = ET.fromstring(document_xml)
    except ET.ParseError:
        return None

    paragraphs = [text for text in _iter_word_like_paragraph_texts(root) if text.strip()]
    object_statistics = _compute_openxml_object_statistics(root, relationships)
    payload = _build_statistics_payload(
        source="openxml_word_like",
        engine="openxml-word-like",
        include_textboxes_footnotes_endnotes=False,
        license_status=None,
        statistics_profile="word_web_approx",
        content_scope="main_document_body",
    )
    payload.update(
        {
            "pages": _to_optional_int((cached_statistics or {}).get("pages")),
            "words": _count_word_words(paragraphs),
            "non_asian_words": _count_non_asian_words(paragraphs),
            "asian_characters": _count_asian_characters(paragraphs),
            "characters": _count_characters(paragraphs, include_spaces=False),
            "characters_with_spaces": _count_characters(paragraphs, include_spaces=True),
            "paragraphs": len(paragraphs),
            "lines": _trusted_cached_line_count(cached_statistics, paragraph_count=len(paragraphs)),
            **object_statistics,
        }
    )
    if payload["pages"] is not None or payload["lines"] is not None:
        payload["statistics_warnings"].append("pages_lines_from_cached_docprops")
    return payload


def _apply_aspose_license(aw_module: Any) -> str:
    global _LICENSE_STATUS

    settings = get_settings()
    license_path = (settings.aspose_words_license_path or "").strip()
    if not license_path:
        return "unlicensed"

    with _LICENSE_LOCK:
        if _LICENSE_STATUS is not None:
            return _LICENSE_STATUS

        path = Path(license_path)
        if not path.is_file():
            _LICENSE_STATUS = "license_missing"
            return _LICENSE_STATUS

        try:
            license_obj = aw_module.License()
            license_obj.set_license(str(path))
        except Exception:
            _LICENSE_STATUS = "license_error"
        else:
            _LICENSE_STATUS = "licensed"
        return _LICENSE_STATUS


def _read_cached_docprops_statistics(raw_bytes: bytes) -> dict[str, Any] | None:
    try:
        with ZipFile(BytesIO(raw_bytes)) as archive:
            return _read_cached_docprops_statistics_from_archive(archive)
    except BadZipFile:
        return None


def _read_cached_docprops_statistics_from_archive(archive: ZipFile) -> dict[str, Any] | None:
    try:
        xml_bytes = archive.read("docProps/app.xml")
    except KeyError:
        return None
    try:
        root = ET.fromstring(xml_bytes)
    except ET.ParseError:
        return None

    payload = _build_statistics_payload(
        source="docprops_cached",
        engine="openxml-docprops",
        include_textboxes_footnotes_endnotes=None,
        license_status=None,
        statistics_profile="docprops_cached",
        content_scope="cached_docprops",
    )
    for element_name, key in STATISTIC_FIELDS.items():
        element = root.find(f"{{{APP_PROPERTIES_NS}}}{element_name}")
        payload[key] = _to_optional_int(element.text if element is not None else None)
    return payload


def _read_part_relationships(
    archive: ZipFile,
    part_name: str,
) -> dict[str, dict[str, str | bool]]:
    normalized_name = part_name.lstrip("/")
    directory = posixpath.dirname(normalized_name)
    basename = posixpath.basename(normalized_name)
    rels_name = posixpath.join(directory, "_rels", f"{basename}.rels")
    try:
        root = ET.fromstring(archive.read(rels_name))
    except (KeyError, ET.ParseError):
        return {}

    relationships: dict[str, dict[str, str | bool]] = {}
    for relationship in root.findall(f"{{{RELATIONSHIPS_NS}}}Relationship"):
        rel_id = relationship.get("Id")
        target = relationship.get("Target")
        if not rel_id or not target:
            continue
        target_mode = relationship.get("TargetMode")
        is_external = target_mode == "External"
        if is_external:
            resolved_target = target
        else:
            resolved_target = posixpath.normpath(posixpath.join(directory, target))
        relationships[rel_id] = {
            "target": resolved_target,
            "type": relationship.get("Type") or "",
            "is_external": is_external,
        }
    return relationships


def _iter_word_like_paragraph_texts(root: ET.Element) -> Iterable[str]:
    body = root.find(f".//{_w_tag('body')}")
    container = body if body is not None else root

    def walk(node: ET.Element) -> Iterable[str]:
        if _local_name(node.tag) == "AlternateContent":
            preferred_branch = _select_preferred_alternate_content_branch(node)
            if preferred_branch is not None:
                yield from walk(preferred_branch)
            return
        if node.tag == _w_tag("txbxContent"):
            return
        if node.tag == _w_tag("p"):
            yield _extract_word_like_paragraph_text(node)
            return
        for child in node:
            yield from walk(child)

    yield from walk(container)


def _extract_word_like_paragraph_text(paragraph: ET.Element) -> str:
    parts: list[str] = []

    def walk(node: ET.Element, *, deleted: bool = False, hidden: bool = False) -> None:
        if _local_name(node.tag) == "AlternateContent":
            preferred_branch = _select_preferred_alternate_content_branch(node)
            if preferred_branch is not None:
                walk(preferred_branch, deleted=deleted, hidden=hidden)
            return
        if node.tag in {_w_tag("txbxContent"), _w_tag("moveFrom")}:
            return
        next_deleted = deleted or node.tag in {_w_tag("del"), _w_tag("delText")}
        next_hidden = hidden or _is_hidden_run(node)
        if next_deleted or next_hidden:
            return
        if node.tag == _w_tag("t"):
            parts.append(node.text or "")
            return
        if node.tag == _w_tag("tab"):
            parts.append("\t")
            return
        if node.tag in {_w_tag("br"), _w_tag("cr")}:
            parts.append("\n")
            return
        if node.tag == _w_tag("noBreakHyphen"):
            parts.append("-")
            return
        if node.tag == _w_tag("sym"):
            symbol_text = _decode_symbol(node)
            if symbol_text:
                parts.append(symbol_text)
            return
        if node.tag == _w_tag("instrText"):
            return
        for child in node:
            walk(child, deleted=next_deleted, hidden=next_hidden)

    walk(paragraph)
    return "".join(parts)


def _compute_openxml_object_statistics(
    root: ET.Element,
    relationships: dict[str, dict[str, str | bool]],
) -> dict[str, int]:
    image_targets: list[str] = []
    inline_image_count = 0
    floating_image_count = 0
    linked_image_count = 0
    chart_count = 0
    smartart_count = 0

    def resolve_relationship(rel_id: str) -> tuple[str, bool]:
        relationship = relationships.get(rel_id)
        if not relationship:
            return rel_id, False
        return str(relationship.get("target") or rel_id), bool(relationship.get("is_external"))

    def add_image(rel_id: str, *, linked: bool, floating: bool) -> None:
        nonlocal inline_image_count, floating_image_count, linked_image_count
        target, is_external = resolve_relationship(rel_id)
        image_targets.append(target)
        if floating:
            floating_image_count += 1
        else:
            inline_image_count += 1
        if linked or is_external:
            linked_image_count += 1

    def walk(node: ET.Element, *, in_inline: bool = False, in_anchor: bool = False) -> None:
        nonlocal chart_count, smartart_count
        if _local_name(node.tag) == "AlternateContent":
            preferred_branch = _select_preferred_alternate_content_branch(node)
            if preferred_branch is not None:
                walk(preferred_branch, in_inline=in_inline, in_anchor=in_anchor)
            return

        next_inline = in_inline or node.tag == _wp_tag("inline")
        next_anchor = in_anchor or node.tag == _wp_tag("anchor")
        if node.tag == _a_tag("blip"):
            rel_id = node.get(_r_tag("embed"))
            linked_rel_id = node.get(_r_tag("link"))
            if rel_id:
                add_image(rel_id, linked=False, floating=next_anchor and not next_inline)
            elif linked_rel_id:
                add_image(linked_rel_id, linked=True, floating=next_anchor and not next_inline)
        elif node.tag == _v_tag("imagedata"):
            rel_id = node.get(_r_tag("id")) or node.get(_r_tag("href"))
            if rel_id:
                add_image(rel_id, linked=bool(node.get(_r_tag("href"))), floating=next_anchor and not next_inline)
        elif node.tag == _c_tag("chart") and node.get(_r_tag("id")):
            chart_count += 1
        elif node.tag == _a_tag("graphicData") and node.get("uri") == DGM_GRAPHIC_URI:
            smartart_count += 1

        for child in node:
            walk(child, in_inline=next_inline, in_anchor=next_anchor)

    walk(root)
    return {
        "image_count": len(image_targets),
        "unique_image_count": len(set(image_targets)),
        "inline_image_count": inline_image_count,
        "floating_image_count": floating_image_count,
        "linked_image_count": linked_image_count,
        "chart_count": chart_count,
        "smartart_count": smartart_count,
    }


def _select_preferred_alternate_content_branch(node: ET.Element) -> ET.Element | None:
    for child in list(node):
        if child.tag == f"{{{MC_NS}}}Choice" and len(child):
            return child
    for child in list(node):
        if child.tag == f"{{{MC_NS}}}Fallback" and len(child):
            return child
    return None


def _is_hidden_run(node: ET.Element) -> bool:
    if node.tag != _w_tag("r"):
        return False
    run_properties = node.find(_w_tag("rPr"))
    if run_properties is None:
        return False
    for tag_name in ("vanish", "webHidden"):
        hidden_property = run_properties.find(_w_tag(tag_name))
        if hidden_property is not None and not _is_false_property(hidden_property.get(_w_tag("val"))):
            return True
    return False


def _is_false_property(value: str | None) -> bool:
    return (value or "").strip().lower() in {"0", "false", "off"}


def _decode_symbol(node: ET.Element) -> str:
    raw_value = node.get(_w_tag("char"))
    if not raw_value:
        return ""
    try:
        return chr(int(raw_value, 16))
    except (TypeError, ValueError):
        return ""


def _iter_ppt_shapes(shapes: Any) -> Iterable[Any]:
    for shape in shapes:
        yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_ppt_shapes(nested_shapes)


def _extract_ppt_shape_texts(shape: Any) -> tuple[list[str], int, int, int]:
    texts: list[str] = []
    paragraph_count = 0
    line_count = 0
    chart_count = 0

    if getattr(shape, "has_text_frame", False):
        text, paragraphs, lines = _extract_ppt_text_frame(shape.text_frame)
        if text:
            texts.append(text)
            paragraph_count += paragraphs
            line_count += lines

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                text, paragraphs, lines = _extract_ppt_text_frame(cell.text_frame)
                if text:
                    texts.append(text)
                    paragraph_count += paragraphs
                    line_count += lines

    if getattr(shape, "has_chart", False):
        chart_count = 1
        chart_text = _extract_ppt_chart_text(shape)
        if chart_text:
            texts.append(chart_text)
            paragraph_count += _count_nonempty_text_lines(chart_text)
            line_count += _count_nonempty_text_lines(chart_text)

    return texts, paragraph_count, line_count, chart_count


def _extract_ppt_text_frame(text_frame: Any) -> tuple[str, int, int]:
    texts: list[str] = []
    paragraph_count = 0
    line_count = 0
    for paragraph in getattr(text_frame, "paragraphs", []) or []:
        text = (paragraph.text or "").strip()
        if not text:
            continue
        texts.append(text)
        paragraph_count += 1
        line_count += _count_nonempty_text_lines(text)

    if not texts:
        text = (getattr(text_frame, "text", "") or "").strip()
        if text:
            texts.append(text)
            paragraph_count = _count_text_blocks(text)
            line_count = _count_nonempty_text_lines(text)
    return "\n".join(texts), paragraph_count, line_count


def _extract_ppt_chart_text(shape: Any) -> str:
    try:
        chart_space = getattr(shape.chart, "_chartSpace", None)
        if chart_space is None:
            return ""
        texts: list[str] = []
        seen: set[str] = set()
        for node in chart_space.iter():
            local_name = _local_name(str(node.tag))
            value = (node.text or "").strip()
            if not value:
                continue
            if local_name == "v":
                try:
                    float(value)
                except ValueError:
                    pass
                else:
                    continue
            elif local_name != "t":
                continue
            if value not in seen:
                seen.add(value)
                texts.append(value)
        return "\n".join(texts)
    except Exception:
        return ""


def _extract_ppt_notes_text(slide: Any) -> tuple[str, int, int]:
    try:
        if not getattr(slide, "has_notes_slide", False):
            return "", 0, 0
        notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
        if notes_frame is None:
            return "", 0, 0
        return _extract_ppt_text_frame(notes_frame)
    except Exception:
        return "", 0, 0


def _count_pptx_image_references(raw_bytes: bytes) -> tuple[int, int]:
    image_count = 0
    linked_image_count = 0
    try:
        with ZipFile(BytesIO(raw_bytes)) as archive:
            for part_name in archive.namelist():
                if not part_name.startswith("ppt/slides/") or not part_name.endswith(".xml"):
                    continue
                try:
                    root = ET.fromstring(archive.read(part_name))
                except (KeyError, ET.ParseError):
                    continue
                for node in root.iter(_a_tag("blip")):
                    embed_id = node.get(_r_tag("embed"))
                    link_id = node.get(_r_tag("link"))
                    if embed_id or link_id:
                        image_count += 1
                    if link_id:
                        linked_image_count += 1
    except BadZipFile:
        return 0, 0
    return image_count, linked_image_count


def _count_tool_word_like_texts(texts: Iterable[str]) -> dict[str, int]:
    asian_characters = 0
    non_asian_words = 0
    characters = 0
    characters_with_spaces = 0

    for text in texts:
        content = text or ""
        characters += sum(1 for char in content if not char.isspace())
        characters_with_spaces += len(content)
        index = 0
        while index < len(content):
            char = content[index]
            if _ppt_cjk_script(char):
                asian_characters += 1
                index += 1
                continue
            if _is_ppt_token_char(char):
                _, index = _consume_ppt_word_like_token(content, index)
                non_asian_words += 1
                continue
            previous_char = content[index - 1] if index > 0 else ""
            next_char = content[index + 1] if index + 1 < len(content) else ""
            if _is_ppt_cjk_punctuation(char, previous_char, next_char):
                asian_characters += 1
            index += 1

    return {
        "words": asian_characters + non_asian_words,
        "non_asian_words": non_asian_words,
        "asian_characters": asian_characters,
        "characters": characters,
        "characters_with_spaces": characters_with_spaces,
    }


def _consume_ppt_word_like_token(text: str, start: int) -> tuple[str, int]:
    index = start
    while index < len(text):
        char = text[index]
        if _is_ppt_token_char(char) or char in {"'", "’", "-", "_", ".", "/"}:
            index += 1
            continue
        break
    return text[start:index], index


def _ppt_cjk_script(char: str) -> str:
    if not char or len(char) != 1:
        return ""
    code = ord(char)
    if (
        0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0xF900 <= code <= 0xFAFF
        or 0x20000 <= code <= 0x2A6DF
        or 0x2A700 <= code <= 0x2B73F
        or 0x2B740 <= code <= 0x2B81F
        or 0x2B820 <= code <= 0x2CEAF
        or 0x2CEB0 <= code <= 0x2EBEF
        or 0x2F800 <= code <= 0x2FA1F
        or 0x30000 <= code <= 0x3134F
    ):
        return "han"
    if 0x3040 <= code <= 0x30FF or 0x31F0 <= code <= 0x31FF or 0xFF66 <= code <= 0xFF9D:
        return "kana"
    if (
        0x1100 <= code <= 0x11FF
        or 0x3130 <= code <= 0x318F
        or 0xA960 <= code <= 0xA97F
        or 0xAC00 <= code <= 0xD7AF
        or 0xD7B0 <= code <= 0xD7FF
    ):
        return "hangul"
    return ""


def _is_ppt_token_char(char: str) -> bool:
    if not char or _ppt_cjk_script(char):
        return False
    return unicodedata.category(char)[0] in {"L", "N"}


def _is_ppt_cjk_punctuation(char: str, previous_char: str = "", next_char: str = "") -> bool:
    if not char or char.isspace() or _is_ppt_token_char(char):
        return False
    code = ord(char)
    if unicodedata.category(char)[0] not in {"P", "S"}:
        return False
    if (
        0x3000 <= code <= 0x303F
        or 0xFE10 <= code <= 0xFE1F
        or 0xFE30 <= code <= 0xFE4F
        or 0xFF00 <= code <= 0xFFEF
    ):
        return True
    return char in {"“", "”", "‘", "’", "—", "–", "…", "·"} and bool(
        _ppt_cjk_script(previous_char) or _ppt_cjk_script(next_char)
    )


def _count_nonempty_text_lines(text: str) -> int:
    return sum(1 for line in re.split(r"\r\n|\r|\n", text or "") if line.strip())


def _count_text_blocks(text: str) -> int:
    content = (text or "").strip()
    if not content:
        return 0
    blocks = [block for block in re.split(r"(?:\r?\n\s*){2,}", content) if block.strip()]
    return len(blocks) if blocks else _count_nonempty_text_lines(content)


def _count_word_words(paragraphs: Iterable[str]) -> int:
    text_parts = list(paragraphs)
    return _count_asian_characters(text_parts) + _count_non_asian_words(text_parts)


def _count_asian_characters(paragraphs: Iterable[str]) -> int:
    return sum(
        1
        for text in paragraphs
        for char in text
        if _is_east_asian_word_count_char(char)
    )


def _count_non_asian_words(paragraphs: Iterable[str]) -> int:
    word_count = 0
    for text in paragraphs:
        non_asian_text = "".join(
            " " if _is_east_asian_word_count_char(char) else char
            for char in text
        )
        word_count += len(NON_ASIAN_WORD_PATTERN.findall(non_asian_text))
    return word_count


def _count_characters(paragraphs: Iterable[str], *, include_spaces: bool) -> int:
    count = 0
    for text in paragraphs:
        for char in text:
            if char in "\r\n":
                continue
            if not include_spaces and char.isspace():
                continue
            count += 1
    return count


def _trusted_cached_line_count(
    cached_statistics: dict[str, Any] | None,
    *,
    paragraph_count: int,
) -> int | None:
    cached_lines = _to_optional_int((cached_statistics or {}).get("lines"))
    if cached_lines is None or cached_lines < paragraph_count:
        return None
    return cached_lines


def _is_east_asian_word_count_char(char: str) -> bool:
    if char.isspace() or char.isascii():
        return False
    return unicodedata.east_asian_width(char) in {"W", "F"}


def _local_name(tag: str) -> str:
    if "}" not in tag:
        return tag
    return tag.rsplit("}", 1)[1]


def _w_tag(local_name: str) -> str:
    return f"{{{WORD_NS}}}{local_name}"


def _r_tag(local_name: str) -> str:
    return f"{{{REL_NS}}}{local_name}"


def _a_tag(local_name: str) -> str:
    return f"{{{A_NS}}}{local_name}"


def _c_tag(local_name: str) -> str:
    return f"{{{C_NS}}}{local_name}"


def _wp_tag(local_name: str) -> str:
    return f"{{{WP_NS}}}{local_name}"


def _v_tag(local_name: str) -> str:
    return f"{{{VML_NS}}}{local_name}"


def _build_statistics_payload(
    *,
    source: str,
    engine: str | None,
    include_textboxes_footnotes_endnotes: bool | None,
    license_status: str | None,
    engine_version: str | None = None,
    statistics_profile: str | None = None,
    content_scope: str | None = None,
    statistics_warnings: Iterable[str] | None = None,
) -> dict[str, Any]:
    payload: dict[str, Any] = {
        "source": source,
        "engine": engine,
        "engine_version": engine_version,
        "license_status": license_status,
        "include_textboxes_footnotes_endnotes": include_textboxes_footnotes_endnotes,
        "statistics_profile": statistics_profile or _default_statistics_profile(source),
        "content_scope": content_scope or _default_content_scope(source),
        "statistics_warnings": list(statistics_warnings or []),
        "match_analysis": None,
    }
    for key in STATISTIC_NUMBER_KEYS:
        payload[key] = None
    return payload


def _coerce_statistics_payload(value: dict[str, Any]) -> dict[str, Any]:
    source = str(value.get("source") or "")
    payload = _build_statistics_payload(
        source=source,
        engine=value.get("engine") if value.get("engine") is None else str(value.get("engine")),
        engine_version=(
            value.get("engine_version")
            if value.get("engine_version") is None
            else str(value.get("engine_version"))
        ),
        include_textboxes_footnotes_endnotes=_to_optional_bool(
            value.get("include_textboxes_footnotes_endnotes")
        ),
        license_status=(
            value.get("license_status")
            if value.get("license_status") is None
            else str(value.get("license_status"))
        ),
        statistics_profile=(
            value.get("statistics_profile")
            if value.get("statistics_profile") is None
            else str(value.get("statistics_profile"))
        ),
        content_scope=(
            value.get("content_scope")
            if value.get("content_scope") is None
            else str(value.get("content_scope"))
        ),
        statistics_warnings=_to_string_list(value.get("statistics_warnings")),
    )
    for key in STATISTIC_NUMBER_KEYS:
        payload[key] = _to_optional_int(value.get(key))
    payload["match_analysis"] = normalize_document_match_analysis(value.get("match_analysis"))
    return payload


def _default_statistics_profile(source: str) -> str:
    if source in {"openxml_word_like", "openxml_computed"}:
        return "word_web_approx"
    if source == "aspose":
        return "aspose_word"
    if source == "libreoffice":
        return "libreoffice"
    if source == "docprops_cached":
        return "docprops_cached"
    return source or "unknown"


def _default_content_scope(source: str) -> str:
    if source in {"openxml_word_like", "openxml_computed"}:
        return "main_document_body"
    if source == "aspose":
        return "aspose_document"
    if source == "libreoffice":
        return "libreoffice_default"
    if source == "docprops_cached":
        return "cached_docprops"
    return "unknown"


def _configured_license_status_hint() -> str:
    license_path = (get_settings().aspose_words_license_path or "").strip()
    if not license_path:
        return "unlicensed"
    return "not_applied"


def _first_int(*values: Any) -> int | None:
    for value in values:
        converted = _to_optional_int(value)
        if converted is not None:
            return converted
    return None


def _to_optional_int(value: Any) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _to_optional_bool(value: Any) -> bool | None:
    if value is None:
        return None
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {"1", "true", "yes", "on"}:
            return True
        if normalized in {"0", "false", "no", "off"}:
            return False
    return bool(value)


def _to_string_list(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(item) for item in value if item is not None]
    if isinstance(value, str) and value.strip():
        return [value]
    return []
