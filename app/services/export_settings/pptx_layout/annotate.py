"""
annotate.py —— 生成"注释版"PPTX：给候选文本框描蓝框、强制全显溢出文字、贴可见数字标签。

改编自参考实现 app/文件/play2.py：移除文件级批处理 demo 与 __main__。
视觉模型据这些蓝框与标签来判断溢出并回传结果。
"""
from __future__ import annotations

import logging

from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def add_tag_label(slide, tag_text: str, abs_left_emu: int, abs_top_emu: int):
    """在注释版 pptx 上贴一个可见的红字黄底数字标签(id)。

    视觉模型用该标签识别每个文本框，并在结构化 JSON 中回传标签号，
    上层再按 标签->uid 精确写回目标框。位置为幻灯片绝对坐标(EMU)。
    """
    w, h = Inches(0.55), Inches(0.28)
    tb = slide.shapes.add_textbox(int(abs_left_emu), int(abs_top_emu), w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_top = 0
    tf.text = tag_text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0)
    try:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(255, 255, 0)
    except Exception:  # noqa: BLE001
        logger.debug("给标签设置黄色填充失败，忽略。", exc_info=True)
    return tb


def draw_border_and_force_visible(shape) -> None:
    """给文本框描蓝框并强制全显溢出文字，不修改任何文本内容与样式。

    1. 通过底层 XML 打开垂直/水平溢出渲染（防止超边文字被截断隐藏）。
    2. 描蓝色边框做可视化标注（100% 保留原本物理宽高与位置）。
    """
    tf = shape.text_frame
    try:
        body_pr = tf._txBody.bodyPr
        body_pr.set("vertOverflow", "overflow")
        body_pr.set("horzOverflow", "overflow")
    except Exception:  # noqa: BLE001
        logger.debug("强制全显 XML 注入失败，忽略。", exc_info=True)

    try:
        line = shape.line
        line.color.rgb = RGBColor(0, 102, 204)
        line.width = Pt(1.5)
    except Exception:  # noqa: BLE001
        logger.debug("描蓝色边框失败，忽略。", exc_info=True)
