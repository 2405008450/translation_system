"""
pipeline.py —— PPTX 版式优化主控流水线（字节进 / 字节出）。

改编自参考实现 app/文件/pipeline.py：
  - 输入/输出为 bytes，全程在临时目录中处理并清理。
  - 渲染或视觉模型不可用时，自动回退到启发式估算（estimate_required_height），
    保证始终产出合法 PPTX；任何未预期异常返回原始字节。

流程：
  1) 采集文本框候选（uid + 位置 + 字号）
  2) 生成蓝框注释版 pptx（供截图）
  3) 逐页截图 -> 视觉模型判断溢出/富余；失败则启发式估算
  4) 按模式把结果写回原始 pptx 的对应文本框
"""
from __future__ import annotations

import logging
import tempfile
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.services.export_settings.pptx_layout.annotate import (
    add_tag_label,
    draw_border_and_force_visible,
)
from app.services.export_settings.pptx_layout.fitter import fit_textbox
from app.services.export_settings.pptx_layout.slide_render import export_slide_to_image
from app.services.export_settings.pptx_layout.textbox_common import (
    collect_layout_records,
    estimate_required_height,
    find_shape_and_xform_by_uid,
    find_shape_by_uid,
    overflow_ratio,
)
from app.services.export_settings.pptx_layout.report_service import persist_pptx_layout_report
from app.services.export_settings.pptx_layout.vlm_client import (
    analyze_page_overflow,
    resolve_vlm_model,
    vlm_provider,
)

logger = logging.getLogger(__name__)

VALID_MODES = {"model_scale", "shrink", "both", "expand"}
DEFAULT_MODE = "model_scale"

_MIN_BOX_HEIGHT = 0.15        # 极小高度框（自动伸缩/装饰）不参与调整
_OVERFLOW_RATIO_GATE = 1.05   # 启发式回退时的溢出阈值


def optimize_pptx_layout(
    pptx_bytes: bytes,
    *,
    mode: str = DEFAULT_MODE,
    model: str | None = None,
    max_expand=None,
    report_context: dict | None = None,
) -> bytes:
    """对 PPTX 字节执行版式优化，返回调整后的字节；失败或无调整时返回原始字节。

    若提供 report_context，则把本次的 AI 判断与实际调整结果落库（失败不影响导出）。
    """
    if mode not in VALID_MODES:
        mode = DEFAULT_MODE

    resolved_provider = vlm_provider()
    resolved_model = resolve_vlm_model(model)

    tmp_dir = Path(tempfile.mkdtemp(prefix="pptx_layout_"))
    input_path = tmp_dir / "input.pptx"
    annotated_path = tmp_dir / "annotated.pptx"
    output_path = tmp_dir / "output.pptx"
    try:
        input_path.write_bytes(pptx_bytes)

        records = collect_layout_records(input_path)
        candidates = [r for r in records if r["height"] >= _MIN_BOX_HEIGHT]
        if not candidates:
            return pptx_bytes

        tag_map = _build_tag_map(candidates)
        _make_annotated(input_path, candidates, annotated_path)
        results, vlm_used = _analyze_all_pages(annotated_path, candidates, tag_map, mode=mode, model=model)
        if not results:
            _safe_persist_report(
                report_context, mode=mode, provider=resolved_provider, model=resolved_model,
                total_candidates=len(candidates), vlm_used=vlm_used, status="no_change", items=[],
            )
            return pptx_bytes

        applied, item_records = _apply_results(
            input_path, results, output_path, candidates, mode=mode, max_expand=max_expand
        )
        _safe_persist_report(
            report_context, mode=mode, provider=resolved_provider, model=resolved_model,
            total_candidates=len(candidates), vlm_used=vlm_used,
            status="completed" if applied else "no_change", items=item_records,
        )
        if applied and output_path.is_file():
            return output_path.read_bytes()
        return pptx_bytes
    except Exception:  # noqa: BLE001
        logger.exception("PPTX 版式优化失败，返回未调整的原始导出内容。")
        _safe_persist_report(
            report_context, mode=mode, provider=resolved_provider, model=resolved_model,
            total_candidates=0, vlm_used=False, status="failed", items=[],
        )
        return pptx_bytes
    finally:
        _cleanup_dir(tmp_dir)


def _safe_persist_report(report_context, **kwargs) -> None:
    if not report_context:
        return
    filename = str((report_context or {}).get("filename") or "")
    persist_pptx_layout_report(report_context, filename=filename, **kwargs)


def _build_tag_map(candidates: list[dict]) -> dict[str, str]:
    """给每个候选框分配页内连番可见标签，返回 tag -> uid。"""
    tag_map: dict[str, str] = {}
    per_slide_counter: dict[int, int] = defaultdict(int)
    for c in candidates:
        per_slide_counter[c["slide_index"]] += 1
        tag = f"{c['slide_index'] + 1}-{per_slide_counter[c['slide_index']]}"
        c["tag"] = tag
        tag_map[tag] = c["uid"]
    return tag_map


def _make_annotated(input_path: Path, candidates: list[dict], output_path: Path) -> None:
    """生成蓝框 + 数字标签的注释版 pptx。"""
    prs = Presentation(str(input_path))
    for c in candidates:
        shape = find_shape_by_uid(prs, c["uid"])
        if shape is None:
            continue
        draw_border_and_force_visible(shape)
        slide = prs.slides[c["slide_index"]]
        add_tag_label(slide, c["tag"], Inches(c["left"]), Inches(c["top"]))
    prs.save(str(output_path))


def _analyze_all_pages(
    annotated_path: Path,
    candidates: list[dict],
    tag_map: dict[str, str],
    *,
    mode: str,
    model: str | None,
) -> tuple[dict[str, dict], bool]:
    """逐页截图 + 视觉复核，得到 (uid -> 结果, 是否用到视觉模型)；视觉不可用时回退启发式估算。"""
    by_slide: dict[int, list[dict]] = defaultdict(list)
    for c in candidates:
        by_slide[c["slide_index"]].append(c)

    results: dict[str, dict] = {}
    vlm_used = False
    tmp_dir = annotated_path.parent
    for si in sorted(by_slide):
        boxes = by_slide[si]
        image_path = tmp_dir / f"slide_{si}.jpg"
        model_boxes = [
            {
                "id": c["tag"],
                "left": c["left"],
                "top": c["top"],
                "width": c["width"],
                "height": c["height"],
                "text": c["text"],
            }
            for c in boxes
        ]

        page_results = None
        if export_slide_to_image(annotated_path, si, image_path):
            page_results = analyze_page_overflow(
                str(image_path), model_boxes, mode=mode, model=model
            )

        if page_results is None:
            # 渲染或视觉模型不可用 -> 启发式估算回退
            results.update(_estimate_fallback(boxes))
            continue

        vlm_used = True
        for tag, box in page_results.items():
            uid = tag_map.get(tag)
            if uid is not None:
                results[uid] = box

    return results, vlm_used


def _estimate_fallback(boxes: list[dict]) -> dict[str, dict]:
    """视觉模型不可用时：对 overflow_ratio 超阈值的框按启发式给出理想框（仅扩高）。"""
    results: dict[str, dict] = {}
    for c in boxes:
        ratio = overflow_ratio(c)
        if ratio > _OVERFLOW_RATIO_GATE:
            est_h = estimate_required_height(c)
            results[c["uid"]] = {
                "left": c["left"],
                "top": c["top"],
                "width": c["width"],
                "height": est_h,
                "overflow_ratio": ratio,
                "_estimate": True,
            }
    return results


def _item_kind(target: dict) -> str:
    if target.get("_estimate"):
        return "estimate"
    if target.get("underflow"):
        return "underflow"
    return "overflow"


def _apply_results(
    input_path: Path,
    results: dict[str, dict],
    output_path: Path,
    candidates: list[dict],
    *,
    mode: str,
    max_expand,
) -> tuple[bool, list[dict]]:
    """把结果写回原始 pptx 对应文本框，返回 (是否有调整, 报告条目列表)。"""
    prs = Presentation(str(input_path))
    rec_by_uid = {c["uid"]: c for c in candidates}
    applied = 0
    items: list[dict] = []
    for uid, target in results.items():
        rec = rec_by_uid.get(uid, {})
        has_geometry = all(k in target for k in ("left", "top", "width", "height"))
        scale = 1.0
        changed = False
        try:
            shape, xform = find_shape_and_xform_by_uid(prs, uid)
            if shape is None:
                continue
            scale = fit_textbox(shape, target, mode=mode, max_expand=max_expand, xform=xform)
            # 判定是否产生实际调整：字号有缩放，或 expand/both 模式移动/放大了框
            geometry_changed = (
                mode in ("expand", "both")
                and has_geometry
                and not target.get("underflow")
            )
            changed = (scale is not None and abs(scale - 1.0) > 1e-6) or geometry_changed
            if changed:
                applied += 1
        except Exception:  # noqa: BLE001
            logger.warning("写回文本框 %s 失败，跳过该框。", uid, exc_info=True)
            continue

        items.append(
            {
                "uid": uid,
                "slide_index": rec.get("slide_index", 0),
                "tag": rec.get("tag", ""),
                "kind": _item_kind(target),
                "source_text": rec.get("text", ""),
                "orig": (rec.get("left"), rec.get("top"), rec.get("width"), rec.get("height")),
                "new": (
                    (target.get("left"), target.get("top"), target.get("width"), target.get("height"))
                    if has_geometry
                    else None
                ),
                "overflow_ratio": target.get("overflow_ratio"),
                "font_scale": scale,
                "applied": changed,
                "reason": target.get("reason", ""),
            }
        )

    if applied:
        prs.save(str(output_path))
    return applied > 0, items


def _cleanup_dir(tmp_dir: Path) -> None:
    try:
        for path in tmp_dir.glob("*"):
            try:
                path.unlink()
            except OSError:
                pass
        tmp_dir.rmdir()
    except OSError:
        logger.debug("清理 PPTX 版式优化临时目录失败：%s", tmp_dir, exc_info=True)
