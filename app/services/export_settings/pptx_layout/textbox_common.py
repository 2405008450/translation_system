"""
textbox_common.py —— PPTX 文本框采集与几何/溢出估算的纯函数集合。

改编自参考实现 app/文件/common.py：
  - 去除硬编码路径与 __main__；print 全部移除（纯函数不产生副作用日志）。
  - 顶层裸导入改为标准库/第三方库导入。
逻辑保持与参考一致，便于对照校验：
  - 文本正规化
  - GROUP(组合)递归贯通的文本框迭代器
  - 唯一 uid 生成（页 + 组合路径 + index）
  - 文本框记录采集 collect_textboxes
  - uid -> shape / (shape, xform) 的重定位
  - 行距解析与所需高度估算（溢出判定与字号缩放共用）
"""
from __future__ import annotations

import math
import re
import statistics

from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def get_autofit_scale(shape) -> float:
    """读取 PowerPoint 缓存的"自动缩小"实效比率(bodyPr/normAutofit@fontScale)。

    无/未设置时返回 1.0。fontScale 为千分率(例 "35000"=35%)。
    """
    try:
        body_pr = shape.text_frame._txBody.bodyPr
        autofit = body_pr.find(qn("a:normAutofit")) if body_pr is not None else None
        if autofit is not None and autofit.get("fontScale"):
            return int(autofit.get("fontScale")) / 100000.0
    except Exception:  # noqa: BLE001
        pass
    return 1.0


def clean_text(text: str | None) -> str:
    """去除全部空白用于模糊匹配。"""
    if not text:
        return ""
    return re.sub(r"\s+", "", text)


def normalize_spaces(text: str | None) -> str:
    """把连续空白折叠为单个空格用于展示。"""
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _group_child_transform(group, xform):
    """把 GROUP 子坐标系->父坐标系的仿射变换与外层 xform 合成后返回。"""
    ax, ay, sx, sy = xform
    grp_sp_pr = group._element.find(qn("p:grpSpPr"))
    xfrm = grp_sp_pr.find(qn("a:xfrm")) if grp_sp_pr is not None else None
    if xfrm is None:
        return xform
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    ch_off, ch_ext = xfrm.find(qn("a:chOff")), xfrm.find(qn("a:chExt"))
    if None in (off, ext, ch_off, ch_ext):
        return xform

    ox, oy = int(off.get("x")), int(off.get("y"))
    ecx, ecy = int(ext.get("cx")), int(ext.get("cy"))
    cox, coy = int(ch_off.get("x")), int(ch_off.get("y"))
    ccx, ccy = int(ch_ext.get("cx")), int(ch_ext.get("cy"))
    scale_x = ecx / ccx if ccx else 1.0
    scale_y = ecy / ccy if ccy else 1.0
    lax, lay = ox - cox * scale_x, oy - coy * scale_y
    return (ax + lax * sx, ay + lay * sy, sx * scale_x, sy * scale_y)


def iter_text_shapes(shapes, slide_index, path_prefix="", xform=(0.0, 0.0, 1.0, 1.0)):
    """递归贯通 GROUP，yield (uid, shape, abs_box, xform)。

    abs_box = (left, top, width, height) 为幻灯片绝对坐标(EMU)。
    """
    ax, ay, sx, sy = xform
    for idx, shape in enumerate(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_xform = _group_child_transform(shape, xform)
            yield from iter_text_shapes(
                shape.shapes, slide_index, f"{path_prefix}g{idx}_", child_xform
            )
        else:
            if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                uid = f"s{slide_index}_{path_prefix}{idx}"
                left = shape.left or 0
                top = shape.top or 0
                width = shape.width or 0
                height = shape.height or 0
                abs_box = (ax + left * sx, ay + top * sy, width * sx, height * sy)
                yield uid, shape, abs_box, xform


def collect_font_sizes(shape) -> list[float]:
    """收集文本框内 run/段落字号(pt)，乘以隐藏自动缩小比率得到实际显示字号。"""
    autofit = get_autofit_scale(shape)
    sizes: list[float] = []
    for para in shape.text_frame.paragraphs:
        if para.font.size:
            sizes.append(round(para.font.size.pt * autofit, 1))
        for run in para.runs:
            if run.font.size:
                sizes.append(round(run.font.size.pt * autofit, 1))
    return sizes


def get_line_spacing_info(shape) -> dict:
    """读取该 shape 的真实行距信息（供 estimate_required_height 使用）。

    返回 {"line_h_mult": 相对倍率, "line_h_abs": 绝对行高(inch) 或 None}
    """
    autofit_reduction = 0.0
    try:
        body_pr = shape.text_frame._txBody.bodyPr
        autofit = body_pr.find(qn("a:normAutofit")) if body_pr is not None else None
        if autofit is not None and autofit.get("lnSpcReduction"):
            autofit_reduction = int(autofit.get("lnSpcReduction")) / 100000.0
    except Exception:  # noqa: BLE001
        pass

    pct_values: list[float] = []
    abs_pt_values: list[float] = []
    try:
        for para in shape.text_frame.paragraphs:
            p_pr = para._p.find(qn("a:pPr"))
            ln_spc = p_pr.find(qn("a:lnSpc")) if p_pr is not None else None
            if ln_spc is None:
                continue
            spc_pct = ln_spc.find(qn("a:spcPct"))
            spc_pts = ln_spc.find(qn("a:spcPts"))
            if spc_pct is not None and spc_pct.get("val"):
                pct_values.append(int(spc_pct.get("val")) / 100000.0)
            elif spc_pts is not None and spc_pts.get("val"):
                abs_pt_values.append(int(spc_pts.get("val")) / 100.0)
    except Exception:  # noqa: BLE001
        pass

    if abs_pt_values:
        pt = statistics.median(abs_pt_values)
        return {"line_h_mult": 1.0, "line_h_abs": (pt / 72.0) * (1 - autofit_reduction)}
    if pct_values:
        pct = statistics.median(pct_values)
        return {"line_h_mult": pct * (1 - autofit_reduction), "line_h_abs": None}
    return {"line_h_mult": 1.0 - autofit_reduction, "line_h_abs": None}


def collect_layout_records(pptx_path) -> list[dict]:
    """从 pptx 文件路径打开演示文稿并采集文本框记录（pipeline 入口用）。"""
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    return collect_textboxes(prs)


def collect_textboxes(prs) -> list[dict]:
    """走查整个演示文稿，返回机器可读的文本框记录列表。"""
    sw, sh = prs.slide_width, prs.slide_height
    records: list[dict] = []
    for si, slide in enumerate(prs.slides):
        for uid, shape, abs_box, _xform in iter_text_shapes(slide.shapes, si):
            left, top, width, height = abs_box
            records.append(
                {
                    "uid": uid,
                    "slide_index": si,
                    "text": shape.text_frame.text,
                    "left": left / Inches(1),
                    "top": top / Inches(1),
                    "width": width / Inches(1),
                    "height": height / Inches(1),
                    "rel_left": (left / sw) * 100,
                    "rel_top": (top / sh) * 100,
                    "rel_width": (width / sw) * 100,
                    "rel_height": (height / sh) * 100,
                    "font_sizes": collect_font_sizes(shape),
                    "line_spacing_info": get_line_spacing_info(shape),
                }
            )
    return records


def find_shape_by_uid(prs, uid):
    """返回与 uid 匹配的 shape（无则 None）。"""
    match = re.match(r"s(\d+)_", uid)
    if not match:
        return None
    si = int(match.group(1))
    slide = prs.slides[si]
    for cur_uid, shape, _abs, _xform in iter_text_shapes(slide.shapes, si):
        if cur_uid == uid:
            return shape
    return None


def find_shape_and_xform_by_uid(prs, uid):
    """返回 (shape, xform)。xform 用于把绝对坐标逆变换回子坐标系再写回。无则 (None, None)。"""
    match = re.match(r"s(\d+)_", uid)
    if not match:
        return None, None
    si = int(match.group(1))
    slide = prs.slides[si]
    for cur_uid, shape, _abs, xform in iter_text_shapes(slide.shapes, si):
        if cur_uid == uid:
            return shape, xform
    return None, None


DEFAULT_FONT_PT = 18.0     # 字号不明时的想定值
LINE_SPACING = 1.2         # 既定"单倍行距"系数（lnSpc 未指定时的基准）
PAD_INCH = 0.3             # 上下内边距概算


def _weighted_len(text: str) -> float:
    """全角(CJK 等)=1.0，半角=0.5 的加权字宽。"""
    w = 0.0
    for ch in text:
        w += 1.0 if ord(ch) > 0x2E7F else 0.5
    return w


def estimate_required_height(rec: dict, font_pt: float | None = None) -> float:
    """由框宽、字号、文本量估算所需高度(inch)。"""
    if font_pt is None:
        sizes = rec.get("font_sizes") or []
        font_pt = statistics.median(sizes) if sizes else DEFAULT_FONT_PT

    char_w_inch = font_pt / 72.0
    if char_w_inch <= 0 or rec["width"] <= 0:
        return rec["height"]
    chars_per_line = max(rec["width"] / char_w_inch, 1.0)

    lines = 0
    for para in rec["text"].split("\n"):
        para = para.rstrip()
        if not para:
            lines += 1
            continue
        lines += max(1, math.ceil(_weighted_len(para) / chars_per_line))

    ls_info = rec.get("line_spacing_info")
    if ls_info and ls_info.get("line_h_abs") is not None:
        line_h_inch = ls_info["line_h_abs"]
    elif ls_info and ls_info.get("line_h_mult") is not None:
        line_h_inch = font_pt * LINE_SPACING * ls_info["line_h_mult"] / 72.0
    else:
        line_h_inch = font_pt * LINE_SPACING / 72.0

    return lines * line_h_inch + PAD_INCH


def overflow_ratio(rec: dict) -> float:
    """所需高度 / 当前高度。>1 疑似溢出。"""
    if rec["height"] <= 0:
        return 0.0
    return estimate_required_height(rec) / rec["height"]
