"""
fitter.py —— 把溢出/富余的文本框收拢：按模式调整框几何并/或缩放字号。

改编自参考实现 app/文件/demo2.py：移除 __main__ 与文件级 demo。
print 改为 logging。fit_textbox 支持四种模式：expand / shrink / both / model_scale。
"""
from __future__ import annotations

import logging
import math
import statistics

from pptx.util import Inches, Pt

from app.services.export_settings.pptx_layout.textbox_common import (
    collect_font_sizes,
    estimate_required_height,
    get_line_spacing_info,
)

logger = logging.getLogger(__name__)

MIN_FONT_SCALE = 0.3          # 字号缩小下限
MAX_FONT_SCALE = 5.0          # 字号放大上限
DEFAULT_INHERIT_PT = 18.0     # 继承字号取不到时的回退基准(pt)
MIN_RATIO_TO_SHRINK = 1.15    # model_scale：低于此的 overflow_ratio 视为误差不缩
UNDERFLOW_FILL_MARGIN = 0.10  # underflow 放大时框内保留余白率
MIN_ENLARGE_SCALE = 1.05      # 低于此的放大视为误差不放大


def _resolve_base_pt(shape) -> float:
    """继承字号的 run 的回退基准字号。"""
    sizes = collect_font_sizes(shape)
    return statistics.median(sizes) if sizes else DEFAULT_INHERIT_PT


def scale_text_box_fonts(shape, scale_factor: float) -> None:
    """把文本框内所有段落/run 的字号按比例缩放（继承字号的 run 用基准回退并显式写回）。"""
    base_pt = _resolve_base_pt(shape)
    before_pts: list[float] = []
    after_pts: list[float] = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            before_pts.append(round(paragraph.font.size.pt, 1))
            new_p = round(max(paragraph.font.size.pt * scale_factor, 1.0), 1)
            paragraph.font.size = Pt(new_p)
            after_pts.append(new_p)
        for run in paragraph.runs:
            cur_pt = run.font.size.pt if run.font.size else base_pt
            before_pts.append(round(cur_pt, 1))
            new_size = round(max(cur_pt * scale_factor, 1.0), 1)
            run.font.size = Pt(new_size)
            after_pts.append(new_size)

    if before_pts:
        logger.debug(
            "字号缩放 %s→%s (scale=%.3f, 共%d处)",
            f"{min(before_pts)}~{max(before_pts)}pt",
            f"{min(after_pts)}~{max(after_pts)}pt",
            scale_factor,
            len(before_pts),
        )


def _font_scale_by_area(box_w: float, box_h: float, req_w: float, req_h: float) -> float:
    """面积比缩小率：s = sqrt(A_box / A_req)。"""
    a_box = max(box_w * box_h, 1e-6)
    a_req = max(req_w * req_h, 1e-6)
    if a_req <= a_box:
        return 1.0
    return max(math.sqrt(a_box / a_req), MIN_FONT_SCALE)


def _font_scale_to_fill(box_w: float, box_h: float, req_w: float, req_h: float) -> float:
    """underflow 放大率：填到框面积的 (1-余白率)，上限 MAX_FONT_SCALE。"""
    a_box = max(box_w * box_h, 1e-6) * (1.0 - UNDERFLOW_FILL_MARGIN)
    a_req = max(req_w * req_h, 1e-6)
    if a_req >= a_box:
        return 1.0
    return min(math.sqrt(a_box / a_req), MAX_FONT_SCALE)


def fit_textbox(shape, target: dict, mode: str = "shrink", max_expand=None,
                xform=(0.0, 0.0, 1.0, 1.0), allow_enlarge: bool = True) -> float:
    """把溢出文本框收拢，返回字号缩放率。

    target = {"left","top","width","height"[,"overflow_ratio"|"underflow"]}（英寸，绝对坐标）。
    mode: expand（移动/缩放框，字号不变）/ shrink（框不动，面积比缩字）/
          both（应用理想框 + 残差缩字）/ model_scale（框不动，按 overflow_ratio 缩字）。
    """
    E = Inches(1)
    ax, ay, sx, sy = xform
    sx = sx or 1.0
    sy = sy or 1.0

    old_l = (ax + (shape.left or 0) * sx) / E
    old_t = (ay + (shape.top or 0) * sy) / E
    old_w = ((shape.width or 0) * sx) / E
    old_h = ((shape.height or 0) * sy) / E

    def _set_abs_geometry(l, t, w, h):
        shape.left = int(round((l * E - ax) / sx))
        shape.top = int(round((t * E - ay) / sy))
        shape.width = int(round((w * E) / sx))
        shape.height = int(round((h * E) / sy))

    # underflow：框不动，仅放大字号
    if target.get("underflow"):
        if not allow_enlarge:
            return 1.0
        rec = {
            "text": shape.text_frame.text,
            "width": old_w,
            "height": old_h,
            "font_sizes": collect_font_sizes(shape),
            "line_spacing_info": get_line_spacing_info(shape),
        }
        needed_h = estimate_required_height(rec)
        scale = _font_scale_to_fill(old_w, old_h, old_w, needed_h)
        if scale >= MIN_ENLARGE_SCALE:
            scale_text_box_fonts(shape, scale)
            return scale
        return 1.0

    if mode == "model_scale":
        ratio = target.get("overflow_ratio")
        if not ratio or ratio <= MIN_RATIO_TO_SHRINK:
            return 1.0
        scale = max(math.sqrt(1.0 / ratio), MIN_FONT_SCALE)
        scale_text_box_fonts(shape, scale)
        return scale

    # 以下模式需要模型给出的理想框几何
    req_l = target.get("left")
    req_t = target.get("top")
    req_w = target.get("width")
    req_h = target.get("height")
    if req_l is None or req_t is None or not req_w or not req_h:
        return 1.0
    if req_w <= 0 or req_h <= 0 or old_w <= 0 or old_h <= 0:
        return 1.0

    if mode == "expand":
        _set_abs_geometry(req_l, req_t, req_w, req_h)
        return 1.0

    if mode == "both":
        app_h = min(req_h, max_expand) if max_expand else req_h
        _set_abs_geometry(req_l, req_t, req_w, app_h)
        rec = {
            "text": shape.text_frame.text,
            "width": req_w,
            "height": app_h,
            "font_sizes": collect_font_sizes(shape),
            "line_spacing_info": get_line_spacing_info(shape),
        }
        needed_h = estimate_required_height(rec)
        scale = _font_scale_by_area(req_w, app_h, req_w, needed_h)
        if scale < 0.999:
            scale_text_box_fonts(shape, scale)
        return scale

    # shrink（默认）：框不动，面积比缩字
    scale = _font_scale_by_area(old_w, old_h, req_w, req_h)
    if scale < 0.999:
        scale_text_box_fonts(shape, scale)
    return scale
