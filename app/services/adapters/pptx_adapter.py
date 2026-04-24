"""
PPTX 适配器模块 - 解析 PowerPoint 文件

Requirements: 9.1, 9.2, 9.3, 9.4, 9.5
"""
from io import BytesIO
from typing import List

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table

from app.services.adapters.base import FormatAdapter
from app.services.adapters.exceptions import ParseError
from app.services.adapters.models import (
    BlockNode,
    DocumentAST,
    NodeType,
    ParseResult,
)
from app.services.adapters.segment_extractor import extract_segments


class PptxAdapter(FormatAdapter):
    """PPTX 文件适配器
    
    使用 python-pptx 解析 PowerPoint 文件，提取幻灯片文本和表格。
    """

    def supported_extensions(self) -> List[str]:
        return [".pptx"]

    def parse(self, raw_bytes: bytes) -> ParseResult:
        """解析 PPTX 文件
        
        Args:
            raw_bytes: 文件字节内容
            
        Returns:
            ParseResult: 解析结果
            
        Raises:
            ParseError: 当文件损坏或无法解析时
        """
        if not raw_bytes:
            return ParseResult(
                ast=DocumentAST(nodes=[], source_format=".pptx"),
                segments=[],
                metadata={},
            )
        
        try:
            prs = Presentation(BytesIO(raw_bytes))
        except Exception as e:
            raise ParseError(
                filename="<unknown>",
                reason=f"无法解析 PPTX 文件: {str(e)}"
            )
        
        nodes = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_nodes = self._parse_slide(slide, slide_num)
            nodes.extend(slide_nodes)
        
        ast = DocumentAST(nodes=nodes, source_format=".pptx")
        segments = extract_segments(ast)
        
        return ParseResult(
            ast=ast,
            segments=segments,
            metadata={"slide_count": len(prs.slides)},
        )

    def _parse_slide(self, slide, slide_num: int) -> List[BlockNode]:
        """解析单个幻灯片
        
        Args:
            slide: python-pptx 幻灯片对象
            slide_num: 幻灯片编号（从 1 开始）
            
        Returns:
            List[BlockNode]: 幻灯片中的块级节点列表
        """
        nodes = []
        
        # 提取形状中的文本
        for shape in slide.shapes:
            shape_nodes = self._parse_shape(shape, slide_num)
            nodes.extend(shape_nodes)
        
        # 提取演讲者备注
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                node = BlockNode(
                    node_type=NodeType.NOTE,
                    text_content=notes_text,
                    metadata={"slide": slide_num, "notes": True},
                )
                nodes.append(node)
        
        return nodes

    def _parse_shape(self, shape: BaseShape, slide_num: int) -> List[BlockNode]:
        """解析形状
        
        Args:
            shape: python-pptx 形状对象
            slide_num: 幻灯片编号
            
        Returns:
            List[BlockNode]: 形状中的块级节点列表
        """
        nodes = []
        
        # 处理表格
        if shape.has_table:
            table_node = self._parse_table(shape.table, slide_num)
            if table_node:
                nodes.append(table_node)
            return nodes
        
        # 处理文本框
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # 判断是否是标题
                    is_title = shape.is_placeholder and hasattr(shape, 'placeholder_format')
                    if is_title:
                        try:
                            ph_type = shape.placeholder_format.type
                            # 1 = TITLE, 2 = CENTER_TITLE
                            is_title = ph_type in (1, 2, 3)
                        except Exception:
                            is_title = False
                    
                    node_type = NodeType.HEADING if is_title else NodeType.PARAGRAPH
                    node = BlockNode(
                        node_type=node_type,
                        text_content=text,
                        metadata={"slide": slide_num},
                    )
                    nodes.append(node)
        
        return nodes

    def _parse_table(self, table: Table, slide_num: int) -> BlockNode | None:
        """解析表格
        
        Args:
            table: python-pptx 表格对象
            slide_num: 幻灯片编号
            
        Returns:
            BlockNode | None: 表格节点
        """
        rows = []
        row_count = len(table.rows)
        col_count = len(table.columns)
        
        for row in table.rows:
            cells = []
            for cell in row.cells:
                text = cell.text.strip()
                cell_node = BlockNode(
                    node_type=NodeType.TABLE_CELL,
                    text_content=text,
                )
                cells.append(cell_node)
            
            row_node = BlockNode(
                node_type=NodeType.TABLE_ROW,
                children=cells,
            )
            rows.append(row_node)
        
        return BlockNode(
            node_type=NodeType.TABLE,
            children=rows,
            metadata={
                "slide": slide_num,
                "rows": row_count,
                "columns": col_count,
            },
        )
