import os
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


def clean_text(text):
    """清理字符串中的空白字符，方便在控制台干净地输出预览"""
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def check_and_print_shape(shape, shape_index, prefix=""):
    """
    检查单个形状是否包含文字，如果有则打印
    prefix 用于标记层级（比如组合内部的元素前面加空格缩进）
    """
    # 判断该元素是否包含文本框
    if shape.has_text_frame:
        raw_text = shape.text_frame.text
        # 过滤掉完全为空白的文本框
        if raw_text.strip():
            # 计算位置（如果是组合内元素，其坐标可能是相对于组合框的相对坐标）
            orig_left = shape.left / Inches(1)
            orig_top = shape.top / Inches(1)
            orig_width = shape.width / Inches(1)
            orig_height = shape.height / Inches(1)

            print(f"{prefix}内部编号: Shape[{shape_index}]" + (f" (类型: {shape.shape_type})" if prefix else ""))
            print(
                f"{prefix}位置大小: Left: {orig_left:.2f}\", Top: {orig_top:.2f}\", Width: {orig_width:.2f}\", Height: {orig_height:.2f}\"")
            print(f"{prefix}文本预览: \"{clean_text(raw_text)[:100]}...\"")
            print(f"{prefix}--------------------------------------------------")
            return 1
    return 0


def scan_shapes_recursively(shapes_collection, prefix=""):
    """
    核心递归函数：穿透多层群组(GroupShape)，找出所有隐藏的文本框
    """
    count = 0
    for shape_index, shape in enumerate(shapes_collection):
        # 🔥 关键点：如果是群组/组合形状，深度递归进去遍历
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{prefix}▶ 发现一个 [组合元素 GroupShape]，正在进入内部扫描...")
            # 递归调用自己，传入组合内的 shapes 集合
            count += scan_shapes_recursively(shape.shapes, prefix + "    ")
        else:
            # 如果是常规形状，直接检查
            count += check_and_print_shape(shape, shape_index, prefix)
    return count


def find_absolutely_all_text_boxes(file_path):
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到文件，请检查路径：{file_path}")
        return

    prs = Presentation(file_path)
    total_text_boxes = 0

    print(f"开始深度无死角扫描文件: {os.path.basename(file_path)}")
    print("=" * 60)

    for slide_index, slide in enumerate(prs.slides):
        print(f"\n📂 [第 {slide_index + 1} 页 Slide]")

        # 调用深度递归函数遍历当前页的所有形状
        found_in_slide = scan_shapes_recursively(slide.shapes, prefix="  ")

        total_text_boxes += found_in_slide
        if found_in_slide == 0:
            print("  （本页未发现任何包含文字的元素）")

    print("\n" + "=" * 60)
    print(
        f"扫描完毕！该 PPT 包含 {len(prs.slides)} 页，共抓取到 {total_text_boxes} 个包含文字的文本框（含组合内隐藏元素）。")


# ================= 配置参数并运行 =================
if __name__ == "__main__":
    target_ppt = r"C:\Users\H\Desktop\word解析和还原\P19-36_含不可编辑_新材料产品介绍（中越双语)_bilingual (1).pptx"

    find_absolutely_all_text_boxes(target_ppt)