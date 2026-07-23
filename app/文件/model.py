import os
import re
import json
import platform
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# 【注意】Windows 系统使用 comtypes 调用本地 PowerPoint 软件导出高清图
# Linux/Mac 系统自动改用 LibreOffice(soffice) 命令行转 PDF，再用 pdf2image 转图片
try:
    import comtypes.client
except ImportError:
    comtypes = None

LIBREOFFICE_BIN = os.environ.get("LIBREOFFICE_BIN", "soffice")
PDF_DPI = int(os.environ.get("PDF_EXPORT_DPI", "200"))


# ==================== 1. 导出 PPT 为高分辨率图片 ====================
def export_slide_to_image(ppt_path, slide_index, output_image_path):
    """
    将指定页码的幻灯片导出为图片。
    Windows: 优先用 comtypes 调用本地 PowerPoint 导出。
    Linux/Mac 或 comtypes 不可用: 自动改用 LibreOffice + pdf2image 方案。
    """
    if platform.system() == "Windows" and comtypes:
        return _export_via_powerpoint(ppt_path, slide_index, output_image_path)
    return _export_via_libreoffice(ppt_path, slide_index, output_image_path)


def _export_via_powerpoint(ppt_path, slide_index, output_image_path):
    print(f"正在调用本地 PowerPoint 导出第 {slide_index + 1} 页为图片...")
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        ppt = powerpoint.Presentations.Open(os.path.abspath(ppt_path))

        # PPT 索引从 1 开始
        slide = ppt.Slides(slide_index + 1)
        # 导出为高分辨率图
        slide.Export(os.path.abspath(output_image_path), "JPG", 1920, 1080)

        ppt.Close()
        powerpoint.Quit()
        print(f"🖼️ 图片导出成功: {output_image_path}")
        return True
    except Exception as e:
        print(f"❌ 导出图片失败: {e}")
        return False


def _export_via_libreoffice(ppt_path, slide_index, output_image_path):
    """
    用 LibreOffice 无头模式把整份 pptx 转成 PDF，再用 pdf2image(依赖 poppler)
    取出指定页码存为图片。要求环境已安装:
      - LibreOffice (提供 soffice 命令，可用 LIBREOFFICE_BIN 环境变量指定路径)
      - pdf2image (pip install pdf2image) + poppler-utils (系统包)
    """
    print(f"正在通过 LibreOffice 导出第 {slide_index + 1} 页为图片...")
    out_dir = os.path.dirname(os.path.abspath(output_image_path)) or "."
    os.makedirs(out_dir, exist_ok=True)
    pdf_name = os.path.splitext(os.path.basename(ppt_path))[0] + ".pdf"
    pdf_path = os.path.join(out_dir, pdf_name)

    try:
        subprocess.run(
            [LIBREOFFICE_BIN, "--headless", "--norestore",
             "--convert-to", "pdf", "--outdir", out_dir, os.path.abspath(ppt_path)],
            check=True, timeout=180,
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        )
    except FileNotFoundError:
        print(f"❌ 未找到 LibreOffice 命令 '{LIBREOFFICE_BIN}'，请安装 libreoffice 或设置 LIBREOFFICE_BIN 环境变量。")
        return False
    except subprocess.CalledProcessError as e:
        print(f"❌ LibreOffice 转 PDF 失败: {e.stderr.decode(errors='ignore')}")
        return False
    except subprocess.TimeoutExpired:
        print("❌ LibreOffice 转 PDF 超时。")
        return False

    if not os.path.exists(pdf_path):
        print(f"❌ 未生成预期的 PDF 文件: {pdf_path}")
        return False

    try:
        from pdf2image import convert_from_path
    except ImportError:
        print("❌ 未安装 pdf2image，请先 pip install pdf2image（并确保系统已装 poppler-utils）。")
        return False

    try:
        pages = convert_from_path(pdf_path, dpi=PDF_DPI,
                                   first_page=slide_index + 1, last_page=slide_index + 1)
        if not pages:
            print(f"❌ PDF 中未找到第 {slide_index + 1} 页。")
            return False
        pages[0].save(output_image_path, "JPEG")
        print(f"🖼️ 图片导出成功: {output_image_path}")
        return True
    except Exception as e:
        print(f"❌ PDF 转图片失败: {e}")
        return False
    finally:
        try:
            os.remove(pdf_path)
        except OSError:
            pass


# ==================== 2. 调用多模态大模型 API ====================
def call_vlm_api(image_path, original_box_info):
    """
    调用多模态大模型（如 Qwen2.5-VL / GPT-4o）
    传入图片和当前文本框的原始坐标，让模型通过视觉判断吐出新的高度。
    """
    # 从环境变量获取 API KEY / BASE_URL / 模型名
    api_key = os.environ.get("OPENAI_API_KEY")
    base_url = os.environ.get("BASE_URL")
    model_name = os.environ.get("VLM_MODEL", "google/gemini-3.1-pro-preview")

    print("🧠 正在将图片与物理坐标发送至多模态视觉模型...")

    # 构造精确的 Prompt，强迫大模型按固定格式（JSON）返回，不要解释
    prompt = f"""
    你是一个精通 PPT 自动排版的设计师模型。
    现在有一张双语 PPT 的截图，图中有一个蓝色线框标记的文本框，因为翻译导致字数发生变化，导致文字超出了文本框底边。

    该文本框的当前【物理尺寸和内容】如下:
    - 文本内容: "{original_box_info['text']}"
    - 当前坐标: Left: {original_box_info['left']:.2f}寸, Top: {original_box_info['top']:.2f}寸, Width: {original_box_info['width']:.2f}寸, Height: {original_box_info['height']:.2f}寸

    请你观察图片中文字实际被渲染出来的最佳位置（文字结束的真实边缘位置）。
    在保持 Left, Top, Width 和字号完全不变的前提下，推算出一个能够完美容纳所有文字、不再发生溢出和重叠的【最佳位置：  Left, Top, Width, Height】。

    请严格只返回以下格式的 JSON 字符串，不要包含 ```json 等任何包裹标记和任何多余的解释文字：
    {{"reason": "分析原因","new_left": 浮点数_单位为英寸,"new_top": 浮点数_单位为英寸,"new_width": 浮点数_单位为英寸,"new_height": 新高度数值_浮点数_单位为英寸}}
    """

    # 💡 这里是标准的 OpenAI / DashScope 多模态标准调用结构
    # 由于是测试环境，我们这里做一个“网络请求和解析”的完整模拟演示：
    if not api_key:
        print("💡 [测试桩提示] 未检测到 API_KEY，激活视觉模拟返回...")
        # 模拟大模型看到图像后做出的黄金比例推理结果
        # 假设原始高度 1.5寸 溢出了，视觉判定扩展到 2.25寸 刚刚好
        mock_vlm_response = '{"reason": "中越双语翻译导致行数增加，原始高度1.50寸不足。视觉捕获文字真实下边沿在3.59寸处，计算得出新高度为 2.25 寸。", "new_height": 2.25}'
        return json.loads(mock_vlm_response)

    # 如果配置了真正的 API KEY，这里会运行真实的网络请求：
    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key,
                        base_url=base_url)

        # 将本地图片转换为 Base64
        import base64
        with open(image_path, "rb") as image_file:
            base64_image = base64.b64encode(image_file.read()).decode('utf-8')

        response = client.chat.completions.create(
            model=model_name,  # 环境变量 VLM_MODEL 配置，默认 qwen2.5-vl-72b
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ],
                }
            ],
            response_format={"type": "json_object"}  # 强迫返回干净的 JSON
        )

        result_text = response.choices[0].message.content
        return json.loads(result_text)
    except Exception as e:
        print(f"❌ API 请求或解析失败: {e}")
        return None


# both 模式专用的几何约束追加段：强调"贴着原位置扩展"，抑制大幅平移
_BOTH_MODE_GEO_CONSTRAINT = """
     ★对 new_left / new_top 有【额外强约束】（both 模式专用，非常重要）：
       · 必须尽量贴近当前 Left/Top，只做"贴着原框位置扩展"式的微调，
         而不是把文本框大幅搬移到页面的其他区域；
       · 默认应保持 new_left≈Left、new_top≈Top（允许的自然浮动幅度约为
         当前 Width/Height 的 30% 以内）；
       · 只有当原位置确实无法容纳（如会与旁边内容重叠、或必须往某个方向
         挪一点才放得下）时，才允许做小范围移动，且移动方向要有明确理由；
       · 主要应通过增大 new_width / new_height 来容纳溢出文字，
         而不是整体挪动文本框的位置；
       · 严禁给出与原位置明显不同区域的坐标。
"""

_DEFAULT_MODE_GEO_CONSTRAINT = """
     不溢出不与相邻元素重叠的最佳文本框位置与大小（浮点，单位英寸）。
"""


# ==================== 2b. 批量调用：一页一张图，一次判断该页所有文本框 ====================
def call_vlm_api_batch(image_path, boxes, fit_mode=None):
    """
    一次性把「一整页的截图」+「该页所有文本框(uid/坐标/文本)」发给多模态模型，
    让模型对每个框判断是否溢出并给出建议新高度。
    调用次数 = 页数（而非框数），既能全框送审又不易触发限流。

    boxes: [{id,left,top,width,height,text}, ...]（同一页, id=图中可见的红黄数字标签）
    fit_mode: 当前 pipeline 使用的 FIT_MODE（"both"/"expand"/"shrink"/"model_scale"）。
              仅 "both" 会切换到专用的、约束更强的几何提示词
              （要求 new_left/new_top 贴近原位置，避免模型把框判断成大幅平移）；
              其余模式沿用原有提示词。
    返回: {id: {"left","top","width","height"}, ...}（英寸, 模型建议的理想框, 字号不变时能容纳全文）
          仅包含模型判定为“溢出、需要调整”的框；模型不可用/失败返回 None，交由上层见积回退。
    """
    api_key = os.environ.get("OPENAI_API_KEY")
    base_url = os.environ.get("BASE_URL")#openai/gpt-5-image-mini  google/gemini-3.6-flash anthropic/claude-sonnet-5
    model_name = os.environ.get("VLM_MODEL", "google/gemini-3.1-pro-preview")#anthropic/claude-fable-5 google/gemini-3.5-flash google/gemini-3.1-pro-preview google/gemini-3.1-flash-image

    # 组装该页所有文本框清单（id 与图中每个框左上角的红黄数字标签一一对应）
    lines = []
    for b in boxes:
        lines.append(
            f'- id="{b["id"]}" | Left:{b["left"]:.2f} Top:{b["top"]:.2f} '
            f'Width:{b["width"]:.2f} Height:{b["height"]:.2f}(寸) | '
            f'文本:"{normalize_boxtext(b["text"])}"'
        )
    box_list_text = "\n".join(lines)

    geo_constraint = _BOTH_MODE_GEO_CONSTRAINT if fit_mode == "both" else _DEFAULT_MODE_GEO_CONSTRAINT

    prompt = f"""
你是一个精通 PPT 自动排版的设计师模型。这是一整页 PPT 的截图。
图中每个文本框都用蓝色线框标出，并在其左上角贴有一个【红字黄底的数字标签(id)】。

【重要判定规则】因为开启了文字强制全显，超出文本框的文字**不会被截断**，而是会溢出到蓝色框外部继续显示。
所以不要用“文字有没有被截断”来判断，而要仔细检查**对比“文字实际渲染的边缘”与“该文本框蓝色线框的边界”**：
理想状态是文字**完全落在蓝框内部并四周留有余白**。只要出现以下任一情况都算溢出 overflow=true：
  · 文字越过蓝框的底边/任一边界（延伸到框外）；
  · 文字虽未越界，但已经**紧贴、碰到或几乎顶到**蓝框的边线（没有余白）。
  · 蓝色框或者蓝色框内边界文字被图片等遮挡或者有重叠部分；
只有当文字四周都在蓝框内、且与边线之间留有明显空隙时，才算 overflow=false。

本页所有待检查文本框如下（id 与图中标签一一对应，单位：英寸）：
{box_list_text}

请你依据【图中标签id】逐个定位对应文本框并仔细分析判断，每个框有三种状态（overflow / underflow / 正常，三者互斥）：

一、若文字越界或贴边/顶边（按上面规则视为溢出），overflow=true、underflow=false，并给出两类信息：
  1) overflow_ratio：文字实际渲染总高度 ÷ 框当前 Height，并按“留约10%余白”在测得值上乘 1.15 左右。
     请按真实溢出程度给值，不要夸大：
       · 文字明显超出框（如占框高约2.3倍）→ 填约 2.5；
       · 文字只是轻微贴边/略微顶到边线 → 填约 1.1~1.2；
     （浮点。只有 overflow=true 时才 > 1.0）
  2) new_left/new_top/new_width/new_height：字号不变时能完美容纳全部文字、不溢出不与相邻元素重叠不被图片遮挡的最佳文本框位置与大小，并考虑文字据文本框边线留约10%余白（即文字不要贴边）
  优先最小改动调整，如尽量在原本框位置往上下或左右空白处调整框，注意不要与相邻文字/图片/其他元素重叠（浮点，单位英寸）。

二、若文字明显偏小、框内距离文字留有【大面积空白】（尤其是文字上下高度距离蓝色框远超正常留白，字号明显可以放大），
    且该框**未被遮挡、周围都有足够空白、放大字号后不会与相邻文字/图片/其他元素重叠**，
    则 underflow=true、overflow=false、overflow_ratio=1.0，
    new_* 一律保持当前值（放大由程序按框大小自动计算，你只需判断"能否安全放大"）。
    ⚠只要放大后可能碰到/压到任何相邻元素、或该框本就贴边/被遮挡，就【不要】判 underflow，
      按正常处理（underflow=false）。宁可保守不放大，也不要造成新的重叠。

三、若文字完全在蓝框内、四周留有正常适中余白（既不贴边也非大片空白），
    则 overflow=false、underflow=false、overflow_ratio=1.0，new_* 保持当前值。
{geo_constraint}
必须为每个 id 都返回一条结果。严格只返回如下 JSON（不要 ```json 包裹、不要多余解释）：
{{"results": [{{"id": "对应标签id", "overflow": true/false, "underflow": true/false, "overflow_ratio": 浮点数_大于1, "new_left": 浮点数, "new_top": 浮点数, "new_width": 浮点数, "new_height": 浮点数, "reason": "简述"}}]}}
"""

    print(f"🧠 [批量] 正在发送整页截图 + {len(boxes)} 个文本框清单至视觉模型...")

    if not api_key:
        print("💡 [测试桩提示] 未检测到 API_KEY，批量分析不可用，交由上层见积回退。")
        return None

    try:
        from openai import OpenAI
        import base64
        client = OpenAI(api_key=api_key, base_url=base_url)
        with open(image_path, "rb") as f:
            base64_image = base64.b64encode(f.read()).decode("utf-8")

        response = client.chat.completions.create(
            model=model_name,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
                ],
            }],
            response_format={"type": "json_object"},
        )
        raw_content = response.choices[0].message.content
        data = json.loads(raw_content)
        results = {}
        for item in data.get("results", []):
            is_over = bool(item.get("overflow"))
            # overflow と underflow は排他。overflow 優先で underflow を無効化。
            is_under = bool(item.get("underflow")) and not is_over
            if not is_over and not is_under:
                continue   # 正常な框は返さない
            entry = {}
            if is_over:
                # 几何字段（供 expand/shrink/both 模式）
                try:
                    entry.update({
                        "left": float(item["new_left"]),
                        "top": float(item["new_top"]),
                        "width": float(item["new_width"]),
                        "height": float(item["new_height"]),
                    })
                except (KeyError, TypeError, ValueError):
                    pass
                # 溢出比例字段（供 model_scale 模式：文字实际高度/框高度）
                try:
                    ratio = float(item["overflow_ratio"])
                    if ratio > 1.0:
                        entry["overflow_ratio"] = ratio
                except (KeyError, TypeError, ValueError):
                    pass
            if is_under:
                # 余白過多で字が小さすぎ→枠は動かさずフォント拡大（幾何は保存しない）
                entry["underflow"] = True
            # 分析原因（仅用于日志打印，非必需字段）
            if item.get("reason"):
                entry["reason"] = item["reason"]
            if entry:   # 有溢出几何/比例 或 underflow 标记即保留
                results[str(item["id"])] = entry
        print(results)
        return results
    except Exception as e:
        print(f"❌ 批量 API 请求或解析失败: {e}")
        if "raw_content" in locals():
            print(f"⚠️ 模型原始返回内容: {raw_content}")
        return None


def normalize_boxtext(text, limit=120):
    """把多行文本压成单行并截断，避免 prompt 过长"""
    t = re.sub(r"\s+", " ", text or "").strip()
    return t[:limit] + ("…" if len(t) > limit else "")


# ==================== 3 & 4. 解析结果并写回 PPT ====================
def process_and_update_ppt(ppt_path, slide_index, target_text, new_height):
    """
    解析大模型的指令，将全新的物理高度重新写入到 PPTX 中，完成完美闭环。
    """
    prs = Presentation(ppt_path)
    found = False

    def clean_match(t):
        return re.sub(r"\s+", "", t)

    target_cleaned = clean_match(target_text)

    # 深度遍历，寻找那个对应的 shape 实施手术
    def update_shape_recursively(shapes):
        nonlocal found
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                update_shape_recursively(shape.shapes)
            else:
                if shape.has_text_frame and clean_match(target_text) in clean_match(shape.text_frame.text):
                    print(f"🎯 正在定位原本的组件，准备执行物理尺寸重构...")
                    print(
                        f"📏 正在调整文本框至 -> Left: {shape.left / Inches(1):.2f}寸, Top: {shape.top / Inches(1):.2f}寸, Width: {shape.width / Inches(1):.2f}寸, Height: {new_height:.2f}寸")

                    # 🔥 实施闭环：物理高度重新写入
                    shape.height = Inches(new_height)
                    found = True

    slide = prs.slides[slide_index]
    update_shape_recursively(slide.shapes)

    if found:
        # 保存为修复版文件
        output_fixed_ppt = ppt_path.replace(".pptx", "_fixed_layout.pptx")
        prs.save(output_fixed_ppt)
        print(f"💾 完美闭环！版面已修正，修复版 PPT 已保存至: {output_fixed_ppt}")
    else:
        print("❌ 未在 PPT 中重新匹配到该文本框，写入失败。")


# ==================== 主控流水线 ====================
def run_pipeline():
    # 参数配置
    input_ppt = r"C:\Users\H\Desktop\word解析和还原\parser_style\pptx\文件\P19-36_含不可编辑_新材料产品介绍（中越双语)_bilingual (1).pptx"
    target_slide_idx = 1  # 对应第 2 页
    temp_image = "slide_snapshot.jpg"

    # 我们先对目标文本框进行定位，抓取它的当前绝对位置，准备发送给大模型
    prs = Presentation(input_ppt)
    slide = prs.slides[target_slide_idx]

    # 模拟我们之前锁定的那个溢出的双语长文本框数据
    box_info = {
        "text": "地面、室内外墙面全场景铺贴；适用于混凝土、砖墙面等基础面，更能胜任低吸水率瓷砖、大型大理石/花岗岩铺贴。",
        "left": 5.05,
        "top": 1.34,
        "width": 7.66,
        "height": 1.50  # 假设它原本高度只有 1.5 寸，装不下字
    }

    print("=== 🛠️ 启动智能排版布局优化闭环流水线 ===")

    # 第 1 步：导出当前页面图片
    export_slide_to_image(input_ppt, target_slide_idx, temp_image)

    # 如果本地没有本地 office 导出，为了让 Demo 能向下跑，我们做个兼容机制
    if not os.path.exists(temp_image):
        print(
            "💡 [测试桩提示] 由于没有渲染出真实图片，请确保你的目录下有一张真实 PPT 的截图，命名为 slide_snapshot.jpg 即可触发真实大模型测试。")
        # 创建一个空文件假装图片存在
        with open(temp_image, "w") as f: f.write("image_data")

    # 第 2 步：将图片和当前文本框位置丢给多模态视觉模型
    vlm_result = call_vlm_api(temp_image, box_info)

    if vlm_result and "new_height" in vlm_result:
        # 第 3 步：解析大模型算出的高度
        recommended_height = float(vlm_result["new_height"])
        print(f"✨ 模型返回分析: {vlm_result['reason']}")
        print(f"📊 模型给出的优化高度数值: {recommended_height} 英寸")

        # 第 4 步：将高度反哺写回 PPT 文件，生成最终无溢出版面
        process_and_update_ppt(input_ppt, target_slide_idx, box_info["text"], recommended_height)
    else:
        print("❌ 闭环流水线中断：未获取到多模态大模型的有效尺寸指令。")

    # 清理临时图片
    if os.path.exists(temp_image):
        os.remove(temp_image)


if __name__ == "__main__":
    run_pipeline()