"""
pipeline.py —— 主控流水线（4スクリプトの串联）

流程:
  1) [slide]  文本框の位置・比率・フォントを採取 (uid付きレコード)
  2) [play2]  溢れ候補に青枠を描画＋文字全显し、注釈版pptxを生成
  3) [model]  注釈版を画像化 → マルチモーダルLLMに画像+座標を投げ → new_height
  4) [demo2]  uidで特定した文本框へ高さ変更＋フォント比率縮放を書き戻し → 最終版pptx

各スクリプトの単体実行機能は維持し、本ファイルは関数を import して束ねるだけ。
"""
import os

from pptx import Presentation

from collections import defaultdict

from pptx.util import Inches

from common import (find_shape_by_uid, find_shape_and_xform_by_uid,
                    overflow_ratio, estimate_required_height,
                    normalize_spaces, clean_text)
from slide import collect_layout
from play2 import draw_border_and_force_visible, add_tag_label
from model import export_slide_to_image, call_vlm_api_batch
from demo2 import fit_textbox

# ==================== 設定（パスはここ1か所に統一） ====================
INPUT_PPT = r"测试布局.pptx"
ANNOTATED_PPT = "pipeline_annotated.pptx"      # 青枠注釈版
FINAL_PPT = "pipeline_final.pptx"              # 最終修正版
IMG_DIR = "debug_images"                       # ページ画像の固定出力フォルダ
TEMP_IMAGE = os.path.join(IMG_DIR, "pipeline_slide_{idx}.jpg")  # 例: debug_images/pipeline_slide_0.jpg
KEEP_IMAGES = True                             # True: 画像を残して目視確認 / False: 処理後に削除
TOP_N_CANDIDATES = 8                           # LLMに投げる溢れ候補の上限

# 判定の主体は「モデル」：一页一张图で該当ページの全文本框をまとめて送審する。
# CANDIDATE_MODE は「どのページ/框をモデルに見せるかの範囲」だけを絞る前フィルタ:
#   "all"  : 全文本框を送審（推奨・全框モデル判定）
#   "auto" : overflow_ratio で粗く絞ってから送審（コスト重視）
#   "text" : TARGET_TEXTS を含む框だけ
#   "uids" : TARGET_UIDS で明示指定
CANDIDATE_MODE = "all"
OVERFLOW_RATIO = 1.05                           # auto前フィルタのしきい値
TARGET_TEXTS = []                               # text モード用
TARGET_UIDS = []                                # uids モード用

# LLM が使えない/失敗した時に、見積り高さで自動修正するか
USE_ESTIMATE_FALLBACK = True

# 自動伸縮(auto-grow)/装飾用の極小高さ框は「溢れない」ので調整対象から除外する
MIN_BOX_HEIGHT = 0.15                            # 寸。これ未満は degenerate として除外

# 溢れの収め方:
#   "expand"      : モデルの4値(new_left/top/width/height)を全適用・字号不变
#   "shrink"      : 元の枠を維持・コード側の面積比でフォント縮小
#   "both"        : モデルの4値を全適用した上で、必要なら面積比でフォント縮小
#   "model_scale" : 枠は一切変えず、モデルの overflow_ratio から scale=sqrt(1/ratio) で字号縮小(推奨)
FIT_MODE = "both"
MAX_EXPAND_INCH = None                           # both で高さ上限を掛けたい時のみ数値(寸)。None=モデル値そのまま

# 文字が小さすぎ＆框に余白が多い(underflow)時、フォントを自動拡大するか。
# True: モデルが「安全に拡大可」と判定した框のみ、框大小に合わせて字号を拡大。
ENABLE_UNDERFLOW_ENLARGE = True


def _select_candidates(records):
    # 極小高さ(auto-grow/装飾)の框は除外：高さが自動伸長する框は溢れない
    records = [r for r in records if r["height"] >= MIN_BOX_HEIGHT]
    if CANDIDATE_MODE == "all":
        return records
    if CANDIDATE_MODE == "uids":
        wanted = set(TARGET_UIDS)
        return [r for r in records if r["uid"] in wanted]
    if CANDIDATE_MODE == "text":
        keys = [clean_text(t) for t in TARGET_TEXTS]
        return [r for r in records
                if any(k and k in clean_text(r["text"]) for k in keys)]
    # auto
    cand = [r for r in records if overflow_ratio(r) > OVERFLOW_RATIO]
    cand.sort(key=overflow_ratio, reverse=True)
    return cand[:TOP_N_CANDIDATES]


def stage1_collect(input_ppt):
    print("\n=== [1/4] 文本框採取 (slide) ===")
    records = collect_layout(input_ppt, do_print=False)
    print(f"含文字文本框: {len(records)} 个")
    candidates = _select_candidates(records)

    # 各候補にページ内連番の可視タグを付与し、tag→uid の対応表を作る
    tag_map = {}          # tag(str) -> uid
    per_slide_counter = defaultdict(int)
    for c in candidates:
        per_slide_counter[c["slide_index"]] += 1
        tag = f"{c['slide_index'] + 1}-{per_slide_counter[c['slide_index']]}"  # 例: 2-3
        c["tag"] = tag
        tag_map[tag] = c["uid"]

    slides = sorted({c["slide_index"] for c in candidates})
    print(f"送审范围(模式={CANDIDATE_MODE}): {len(candidates)} 个框, 分布在 {len(slides)} 页")
    return records, candidates, tag_map


def stage2_annotate(input_ppt, candidates, output_ppt):
    """青枠＋文字全显＋可視タグを焼き込んだ注釈版pptxを生成（画像化はこの版を使う）"""
    print("\n=== [2/4] 青枠注釈＋文字全显＋番号タグ (play2) ===")
    prs = Presentation(input_ppt)
    marked = 0
    for c in candidates:
        shape = find_shape_by_uid(prs, c["uid"])
        if shape is None:
            continue
        draw_border_and_force_visible(shape)
        # スライド絶対座標(inch→EMU)にタグを焼き込む
        slide = prs.slides[c["slide_index"]]
        add_tag_label(slide, c["tag"], Inches(c["left"]), Inches(c["top"]))
        marked += 1
    prs.save(output_ppt)
    print(f"注釈済み {marked} 个(含番号タグ) -> {output_ppt}")
    return output_ppt


def _estimate_fallback(boxes):
    """
    モデル不可用時: overflow_ratio 超の框だけ、理想框(target)を見積る。
    位置・幅は現状維持、高さのみ見積り高さに拡張した box を返す。
    返り値: {uid: {"left","top","width","height"}}
    """
    results = {}
    if not USE_ESTIMATE_FALLBACK:
        return results
    for c in boxes:
        if overflow_ratio(c) > OVERFLOW_RATIO:
            est = estimate_required_height(c)
            results[c["uid"]] = {
                "left": c["left"], "top": c["top"],
                "width": c["width"], "height": est,
                "overflow_ratio": overflow_ratio(c),   # model_scale 用の回退比率
            }
            print(f"   ↩ [{c['uid']}] 见积回退 H={c['height']:.2f}→{est:.2f}寸 "
                  f"(ratio≈{overflow_ratio(c):.2f})")
    return results


def stage3_analyze(annotated_ppt, candidates, tag_map):
    print("\n=== [3/4] 画像化＋LLM批量分析 (model, 一页一图/全框送审) ===")
    # ページごとに框をまとめる
    by_slide = defaultdict(list)
    for c in candidates:
        by_slide[c["slide_index"]].append(c)

    results = {}  # uid -> {"left","top","width","height"}(理想框, 英寸)
    for si in sorted(by_slide):
        boxes = by_slide[si]
        img = TEMP_IMAGE.format(idx=si)
        print(f"\n📄 第 {si + 1} 页: {len(boxes)} 个文本框一次性送审")

        # モデルには「可視タグid + 絶対座標 + 文本」を渡す
        model_boxes = [{
            "id": c["tag"],
            "left": c["left"], "top": c["top"],
            "width": c["width"], "height": c["height"],
            "text": c["text"],
        } for c in boxes]

        tag_to_orig = {c["tag"]: c for c in boxes}   # tag -> 原始框记录(打印用)

        ok = export_slide_to_image(annotated_ppt, si, img) if not os.path.exists(img) else True

        page_results = None
        if ok and os.path.exists(img):
            # 発送前に、モデルが実際に読み取る画像パスを表示（目視確認用）
            print(f"   🖼️ 模型将读取的图片(已保存): {os.path.abspath(img)}")
            page_results = call_vlm_api_batch(img, model_boxes, fit_mode=FIT_MODE)

        if page_results is None:
            # モデル不可用（画像化失敗 or API失敗）→ 見積り回退（uidキー）
            print("   ⚠ 模型不可用，改用见积高度回退。")
            results.update(_estimate_fallback(boxes))
        else:
            if not page_results:
                print("   ✅ 模型判定本页无溢出框。")
            for tag, box in page_results.items():
                uid = tag_map.get(tag)
                if uid is None:
                    print(f"   ⚠ 模型返回未知标签 id={tag}，忽略")
                    continue
                reason = box.get("reason", "")
                ratio = box.get("overflow_ratio")
                ratio_s = f"ratio={ratio:.2f} " if ratio else ""
                orig = tag_to_orig.get(tag)
                orig_s = (f"原本框：L{orig['left']:.2f} T{orig['top']:.2f} "
                           f"W{orig['width']:.2f} H{orig['height']:.2f}寸  "
                           if orig else "")
                reason_s = f"{reason} " if reason else ""
                if box.get("underflow"):
                    # 空间富余：只放大字号，框不动
                    print(f"   🔎 [tag {tag} → {uid}] 空间富余(underflow) → "
                          f"{reason_s}{orig_s}将放大字号")
                else:
                    geo_s = (f"理想框 L{box['left']:.2f} T{box['top']:.2f} "
                             f"W{box['width']:.2f} H{box['height']:.2f}寸"
                             if "left" in box else "(无几何)")
                    print(f"   ✨ [tag {tag} → {uid}] 溢出 → {reason_s}{ratio_s}{orig_s}{geo_s}")
                results[uid] = box

    return results


def stage4_apply(input_ppt, results, output_ppt):
    print("\n=== [4/4] 书き戻し＋フォント縮放 (demo2) ===")
    prs = Presentation(input_ppt)
    applied = 0
    for uid, target in results.items():
        shape, xform = find_shape_and_xform_by_uid(prs, uid)
        if shape is None:
            print(f"❌ [{uid}] 未重新定位到 shape，跳过")
            continue
        print(f"🎯 [{uid}] 调整中...")
        # xform を渡し、書き戻し時に絶対座標→子座標へ逆変換（組合内のズレ防止）
        fit_textbox(shape, target, mode=FIT_MODE, max_expand=MAX_EXPAND_INCH, xform=xform,
                    allow_enlarge=ENABLE_UNDERFLOW_ENLARGE)
        applied += 1
    if applied:
        prs.save(output_ppt)
        print(f"✅ 完成 {applied} 个调整 -> {output_ppt}")
    else:
        print("⚠ 无可应用的调整。")


def cleanup():
    if KEEP_IMAGES:
        print(f"🖼️ 页面图片已保留在: {os.path.abspath(IMG_DIR)}")
        return
    if not os.path.isdir(IMG_DIR):
        return
    for f in os.listdir(IMG_DIR):
        if f.startswith("pipeline_slide_") and f.endswith(".jpg"):
            try:
                os.remove(os.path.join(IMG_DIR, f))
            except OSError:
                pass


def run():
    if not os.path.exists(INPUT_PPT):
        print(f"❌ 找不到输入文件: {INPUT_PPT}")
        return

    print("=== 🛠️ 智能排版布局优化闭环流水线 ===")
    os.makedirs(IMG_DIR, exist_ok=True)   # 画像出力フォルダを用意
    records, candidates, tag_map = stage1_collect(INPUT_PPT)
    if not candidates:
        print("未发现溢出候补，流程结束。")
        return

    annotated = stage2_annotate(INPUT_PPT, candidates, ANNOTATED_PPT)
    results = stage3_analyze(annotated, candidates, tag_map)
    if results:
        # 原文ベースに書き戻す（青枠なしの綺麗な最終版）
        stage4_apply(INPUT_PPT, results, FINAL_PPT)
    else:
        print("❌ 未获取到任何有效高度，闭环中断。")

    cleanup()


if __name__ == "__main__":
    run()
