"""
font_report.py —— 字体/字号核对报告
遍历 PPT 全部文本框(含嵌套组合)，打印每个框的位置、文字预览、
以及每一段/每个 run 的字体名与字号——同时给出『名义字号』和
『真实显示字号』(=名义值 × PowerPoint 隐藏自动缩放 normAutofit@fontScale)。

用途: 核对 pipeline 处理前/后的 pptx，字号是否与 PowerPoint 里
肉眼看到的一致（隐藏 autofit 存在时，名义值和真实显示值会不同）。
"""
import os
import sys
from pptx import Presentation
from pptx.util import Inches

from common import iter_text_shapes, get_autofit_scale, normalize_spaces

INPUT_PPT = r"测试布局.pptx"   # 默认核对文件，可用命令行参数覆盖


def _fmt_size(pt):
    return f"{pt:.1f}pt" if pt is not None else "继承(未显式设置)"


def report(ppt_path):
    if not os.path.exists(ppt_path):
        print(f"❌ 找不到文件: {ppt_path}")
        return

    prs = Presentation(ppt_path)
    print("=" * 70)
    print(f"📄 字体/字号核对报告: {os.path.basename(ppt_path)}")
    print("=" * 70)

    total = 0
    for si, slide in enumerate(prs.slides):
        for uid, shape, abs_box, _xform in iter_text_shapes(slide.shapes, si):
            total += 1
            l, t, w, h = [v / Inches(1) for v in abs_box]
            autofit = get_autofit_scale(shape)
            text_preview = normalize_spaces(shape.text_frame.text)[:50]

            print(f"\n[{uid}] 第{si + 1}页  L{l:.2f} T{t:.2f} W{w:.2f} H{h:.2f}寸")
            print(f"  文本预览: \"{text_preview}\"")
            if autofit != 1.0:
                print(f"  ⚠ 检测到 PowerPoint 隐藏自动缩放(normAutofit): {autofit:.3f}  "
                      f"→ 名义字号需 ×{autofit:.3f} 才是肉眼真实看到的大小")

            for pi, para in enumerate(shape.text_frame.paragraphs):
                para_text = para.text.strip()
                if not para_text and not para.runs:
                    continue

                # 段落级字号(若段落本身设置了 font.size)
                if para.font.size:
                    nominal = para.font.size.pt
                    real = nominal * autofit
                    print(f"  段落[{pi}] 段落级字号: 名义={_fmt_size(nominal)} "
                          f"真实显示≈{real:.1f}pt")

                for ri, run in enumerate(para.runs):
                    font_name = run.font.name or "继承(未显式设置)"
                    if run.font.size:
                        nominal = run.font.size.pt
                        real = nominal * autofit
                        size_s = f"名义={nominal:.1f}pt 真实显示≈{real:.1f}pt"
                    else:
                        size_s = "字号=继承母版/占位符(未显式设置)"
                    run_text = run.text.strip()[:30]
                    print(f"  段落[{pi}] run[{ri}] 字体={font_name}  {size_s}  "
                          f"文本=\"{run_text}\"")

    print("\n" + "=" * 70)
    print(f"✅ 共扫描到 {total} 个含文字文本框(含组合内)。")
    print("=" * 70)


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else INPUT_PPT
    report(path)
