"""
common.py —— パイプライン共通土台
- 全スクリプト(play2 / slide / model / demo2)で使う共通ロジックを集約:
  1. テキスト正規化
  2. GROUP(組合)を再帰貫通する文本框イテレータ
  3. 一意な uid の生成（ページ + 組合パス + index）
  4. 文本框レコードの採取 collect_textboxes
  5. uid から shape を再取得する find_shape_by_uid
これにより「slide が採取 → model が分析 → demo2 が書き戻し」で
同じ文本框を uid で確実に指し示せる（特定ズレ・GROUP漏れを防止）。
"""
import re
import math
import statistics
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


def get_autofit_scale(shape):
    """
    読み取り: PowerPoint がキャッシュしている『框に合わせて自動縮小』の実効比率
    (bodyPr の a:normAutofit@fontScale)。
    PowerPoint は run の名義字号(run.font.size)の【上に更に】この比率を掛けて
    実際に表示する。python-pptx で読める run.font.size は名義値のみで、
    この隠れた縮小は含まれない。normAutofit が無い/fontScale未設定なら 1.0(隠れ縮小なし)。
    fontScale は千分率(例: "35000"=35%)。
    """
    try:
        bodyPr = shape.text_frame._txBody.bodyPr
        autofit = bodyPr.find(qn('a:normAutofit')) if bodyPr is not None else None
        if autofit is not None and autofit.get('fontScale'):
            return int(autofit.get('fontScale')) / 100000.0
    except Exception:
        pass
    return 1.0


def clean_text(text):
    """空白・改行を全除去して曖昧一致用に正規化"""
    if not text:
        return ""
    return re.sub(r"\s+", "", text)


def normalize_spaces(text):
    """連続空白を1つに畳んで表示用に整形"""
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def _group_child_transform(group, xform):
    """
    GROUP の子座標系→親座標系の affine を、外側 xform(=(ax,ay,sx,sy)) と合成して返す。
    子点 child を絶対座標へ: abs = a + child * s。
    """
    ax, ay, sx, sy = xform
    grpSpPr = group._element.find(qn('p:grpSpPr'))
    xfrm = grpSpPr.find(qn('a:xfrm')) if grpSpPr is not None else None
    if xfrm is None:
        return xform
    off, ext = xfrm.find(qn('a:off')), xfrm.find(qn('a:ext'))
    chOff, chExt = xfrm.find(qn('a:chOff')), xfrm.find(qn('a:chExt'))
    if None in (off, ext, chOff, chExt):
        return xform

    ox, oy = int(off.get('x')), int(off.get('y'))
    ecx, ecy = int(ext.get('cx')), int(ext.get('cy'))
    cox, coy = int(chOff.get('x')), int(chOff.get('y'))
    ccx, ccy = int(chExt.get('cx')), int(chExt.get('cy'))
    scale_x = ecx / ccx if ccx else 1.0
    scale_y = ecy / ccy if ccy else 1.0
    # ローカル affine: parent = (off - chOff*scale) + child*scale
    lax, lay = ox - cox * scale_x, oy - coy * scale_y
    # 外側 xform と合成
    return (ax + lax * sx, ay + lay * sy, sx * scale_x, sy * scale_y)


def iter_text_shapes(shapes, slide_index, path_prefix="", xform=(0.0, 0.0, 1.0, 1.0)):
    """
    GROUP を再帰貫通し、文字を持つ文本框だけを (uid, shape, abs_box) で yield する。
    abs_box = (left, top, width, height) はスライド絶対座標(EMU)。
    uid 例: s1_3 (2ページ目のindex3) / s1_g0_2 (組合0の中のindex2)
    """
    ax, ay, sx, sy = xform
    for idx, shape in enumerate(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            child_xform = _group_child_transform(shape, xform)
            yield from iter_text_shapes(shape.shapes, slide_index, f"{path_prefix}g{idx}_", child_xform)
        else:
            if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                uid = f"s{slide_index}_{path_prefix}{idx}"
                left = (shape.left or 0)
                top = (shape.top or 0)
                width = (shape.width or 0)
                height = (shape.height or 0)
                abs_box = (ax + left * sx, ay + top * sy, width * sx, height * sy)
                # xform も一緒に返す：書き戻し時に絶対座標→子座標へ逆変換するため
                yield uid, shape, abs_box, xform


def collect_font_sizes(shape):
    """
    文本框内の run/段落フォントサイズ(pt)を収集（縮放計算の参考）。
    ★ PowerPoint の隠れ自動縮小(normAutofit@fontScale)を乗算し、
      『実際に画面へ表示される字号』を返す（名義値のみだと過大評価してしまう）。
    """
    autofit = get_autofit_scale(shape)
    sizes = []
    for para in shape.text_frame.paragraphs:
        if para.font.size:
            sizes.append(round(para.font.size.pt * autofit, 1))
        for run in para.runs:
            if run.font.size:
                sizes.append(round(run.font.size.pt * autofit, 1))
    return sizes


def collect_textboxes(prs):
    """
    プレゼン全体を走査し、共通データ構造のレコード list を返す。
    slide.py の印刷機能とは独立に、機械可読なデータを供給する。
    """
    sw, sh = prs.slide_width, prs.slide_height
    records = []
    for si, slide in enumerate(prs.slides):
        for uid, shape, abs_box, _xform in iter_text_shapes(slide.shapes, si):
            left, top, width, height = abs_box  # スライド絶対座標(EMU, group変換済み)
            records.append({
                "uid": uid,
                "slide_index": si,
                "text": shape.text_frame.text,
                # 絶対サイズ(inch)
                "left": left / Inches(1),
                "top": top / Inches(1),
                "width": width / Inches(1),
                "height": height / Inches(1),
                # 版面比率(%)
                "rel_left": (left / sw) * 100,
                "rel_top": (top / sh) * 100,
                "rel_width": (width / sw) * 100,
                "rel_height": (height / sh) * 100,
                "font_sizes": collect_font_sizes(shape),
                "line_spacing_info": get_line_spacing_info(shape),
            })
    return records


def find_shape_by_uid(prs, uid):
    """uid に一致する shape を返す（無ければ None）"""
    si = int(re.match(r"s(\d+)_", uid).group(1))
    slide = prs.slides[si]
    for cur_uid, shape, _abs, _xform in iter_text_shapes(slide.shapes, si):
        if cur_uid == uid:
            return shape
    return None


def find_shape_and_xform_by_uid(prs, uid):
    """
    uid に一致する shape と、その子座標系→スライド絶対座標への affine 変換
    xform=(ax, ay, sx, sy)（ax/ay は EMU, sx/sy は無次元）を返す。無ければ (None, None)。
    絶対座標(inch)を書き戻す際に、以下で子座標(EMU)へ逆変換して使う:
        child_left_emu  = (left_inch  * 914400 - ax) / sx
        child_top_emu   = (top_inch   * 914400 - ay) / sy
        child_width_emu = (width_inch * 914400) / sx
        child_height_emu= (height_inch* 914400) / sy
    組合の中の形状でも、絶対座標指定した位置へ正しく落とせる（二重変換ズレを防止）。
    """
    si = int(re.match(r"s(\d+)_", uid).group(1))
    slide = prs.slides[si]
    for cur_uid, shape, _abs, xform in iter_text_shapes(slide.shapes, si):
        if cur_uid == uid:
            return shape, xform
    return None, None


DEFAULT_FONT_PT = 18.0     # フォント不明時の想定字号
LINE_SPACING = 1.2         # 既定『単一行間』の係数（lnSpc未指定時の基準値）
PAD_INCH = 0.2             # 上下内边距の概算


def get_line_spacing_info(shape):
    """
    読み取り: 该 shape の真実行間情報（estimate_required_height の
    ハードコード LINE_SPACING(1.2固定) を置き換えるため）。

    OOXML の a:lnSpc/a:spcPct@val は『既定の単一行間(≈フォントサイズ×LINE_SPACING)』
    に対する【相対倍率】であり、170%指定なら実効行高は 1.2×1.70 になる
    （170%そのものを1.2の代わりに使うのは誤り＝過小評価の原因）。
    a:lnSpc/a:spcPts@val は【絶対pt指定】で、フォントサイズと無関係に固定の行高になる。
    さらに normAutofit@lnSpcReduction があれば、PowerPoint が自動で詰めた分として
    どちらにも追加で乗算する。

    戻り値: {"line_h_mult": 相対倍率(spcPct指定時のpct、未指定なら1.0。既に autofit
             縮小を含む), "line_h_abs": 絶対行高(inch)。spcPts指定時のみ値、無ければ None}
    """
    autofit_reduction = 0.0
    try:
        bodyPr = shape.text_frame._txBody.bodyPr
        autofit = bodyPr.find(qn('a:normAutofit')) if bodyPr is not None else None
        if autofit is not None and autofit.get('lnSpcReduction'):
            autofit_reduction = int(autofit.get('lnSpcReduction')) / 100000.0
    except Exception:
        pass

    pct_values, abs_pt_values = [], []
    try:
        for para in shape.text_frame.paragraphs:
            pPr = para._p.find(qn('a:pPr'))
            lnSpc = pPr.find(qn('a:lnSpc')) if pPr is not None else None
            if lnSpc is None:
                continue
            spcPct = lnSpc.find(qn('a:spcPct'))
            spcPts = lnSpc.find(qn('a:spcPts'))
            if spcPct is not None and spcPct.get('val'):
                pct_values.append(int(spcPct.get('val')) / 100000.0)
            elif spcPts is not None and spcPts.get('val'):
                abs_pt_values.append(int(spcPts.get('val')) / 100.0)
    except Exception:
        pass

    if abs_pt_values:
        pt = statistics.median(abs_pt_values)
        return {"line_h_mult": 1.0, "line_h_abs": (pt / 72.0) * (1 - autofit_reduction)}
    if pct_values:
        pct = statistics.median(pct_values)
        return {"line_h_mult": pct * (1 - autofit_reduction), "line_h_abs": None}
    return {"line_h_mult": 1.0 - autofit_reduction, "line_h_abs": None}


def _weighted_len(text):
    """全角(CJK等)=1.0, 半角=0.5 として文字幅を重み付け"""
    w = 0.0
    for ch in text:
        w += 1.0 if ord(ch) > 0x2E7F else 0.5
    return w


def estimate_required_height(rec, font_pt=None):
    """
    枠幅・字号・テキスト量から必要高さ(inch)を概算する。
    - 1全角文字幅 ≈ font_pt/72 inch
    - 1行に入る全角文字数 ≈ width / 文字幅
    - 段落内の改行と明示改行(\n)を合算して行数を出す
    - 行高は rec["line_spacing_info"](= get_line_spacing_info の戻り値)があれば
      その真実行距(a:lnSpc + normAutofit@lnSpcReduction)を使う。無ければ既定 LINE_SPACING。
      ★以前は行距を LINE_SPACING=1.2 固定にしていたため、170%等の広い行距を
        指定した框で必要高さを過小評価していた（＝溢れ誤判定/縮放不足の原因）。
    面積ベースより遥かに実態に近く、溢れ判定と縮放高さの両方に使える。
    """
    if font_pt is None:
        sizes = rec.get("font_sizes") or []
        font_pt = statistics.median(sizes) if sizes else DEFAULT_FONT_PT

    char_w_inch = font_pt / 72.0
    if char_w_inch <= 0 or rec["width"] <= 0:
        return rec["height"]
    chars_per_line = max(rec["width"] / char_w_inch, 1.0)

    lines = 0
    for para in rec["text"].split("\n"):
        para = para.rstrip()
        if not para:
            lines += 1
            continue
        lines += max(1, math.ceil(_weighted_len(para) / chars_per_line))

    ls_info = rec.get("line_spacing_info")
    if ls_info and ls_info.get("line_h_abs") is not None:
        line_h_inch = ls_info["line_h_abs"]                                  # 絶対pt指定(spcPts)
    elif ls_info and ls_info.get("line_h_mult") is not None:
        line_h_inch = font_pt * LINE_SPACING * ls_info["line_h_mult"] / 72.0  # 相対倍率(spcPct, 既定基準×相対%)
    else:
        line_h_inch = font_pt * LINE_SPACING / 72.0                          # 情報無し→既定値

    return lines * line_h_inch + PAD_INCH


def overflow_ratio(rec):
    """必要高さ / 現在高さ。>1 なら溢れ疑い。"""
    if rec["height"] <= 0:
        return 0.0  # 高さ≒0の微小/装飾boxはノイズとして除外
    return estimate_required_height(rec) / rec["height"]
