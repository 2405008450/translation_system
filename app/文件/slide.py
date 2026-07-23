import os
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from common import collect_textboxes, normalize_spaces


def clean_text(text):
    if not text:
        return ""
    return re.sub(r"\s+", " ", text).strip()


def collect_layout(file_path, do_print=True):
    """
    パイプライン向け採取関数。
    共通データ構造(records)を返しつつ、必要なら人間可読なレポートも印刷する。
    ※ 既存の analyze_pptx_layout / process_shapes_relative_position は変更しない。
    """
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到文件：{file_path}")
        return []

    prs = Presentation(file_path)
    records = collect_textboxes(prs)

    if do_print:
        page_w = prs.slide_width / Inches(1)
        page_h = prs.slide_height / Inches(1)
        print("==================================================")
        print("📊 PPT 版面分析 (collect_layout)")
        print(f"   文件名: {os.path.basename(file_path)}")
        print(f"   版面物理大小: 宽 {page_w:.2f} 英寸, 高 {page_h:.2f} 英寸")
        print(f"   采集到含文字文本框: {len(records)} 个")
        print("==================================================")
        for r in records:
            print(f"[{r['uid']}] 绝对: L{r['left']:.2f} T{r['top']:.2f} "
                  f"W{r['width']:.2f} H{r['height']:.2f} 寸 | "
                  f"相对: L{r['rel_left']:.1f}% T{r['rel_top']:.1f}% "
                  f"W{r['rel_width']:.1f}% H{r['rel_height']:.1f}%")
            print(f"       文本: \"{normalize_spaces(r['text'])[:60]}...\"")

    return records


def process_shapes_relative_position(shapes_collection, slide_width, slide_height, slide_index, prefix=""):
    """
    深度递归遍历：计算并打印文本框在 PPT 版面中的绝对大小与相对比例位置
    """
    for shape_index, shape in enumerate(shapes_collection):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"{prefix}▶ [组合元素] 正在进入内部计算相对位置...")
            process_shapes_relative_position(shape.shapes, slide_width, slide_height, slide_index, prefix + "    ")
        else:
            if shape.has_text_frame and shape.text_frame.text.strip():
                # 1. 获取文本框的绝对物理尺寸（英寸）
                abs_left = shape.left / Inches(1)
                abs_top = shape.top / Inches(1)
                abs_width = shape.width / Inches(1)
                abs_height = shape.height / Inches(1)

                # 2. 🔥核心：计算在当前版面中的【相对位置与比例】（百分比）
                # shape.left / slide_width 得到 0~1 之间的比例
                rel_left = (shape.left / slide_width) * 100
                rel_top = (shape.top / slide_height) * 100
                rel_width = (shape.width / slide_width) * 100
                rel_height = (shape.height / slide_height) * 100

                # 3. 打印模拟输出
                raw_text = shape.text_frame.text
                print(f"{prefix}内部编号: Shape[{shape_index}]")
                print(
                    f"{prefix}绝对尺寸: Left: {abs_left:.2f}\", Top: {abs_top:.2f}\", Width: {abs_width:.2f}\", Height: {abs_height:.2f}\"")
                print(
                    f"{prefix}相对版面: Left: {rel_left:.1f}%, Top: {rel_top:.1f}%, Width: {rel_width:.1f}%, Height: {rel_height:.1f}%")
                print(f"{prefix}文本内容: \"{clean_text(raw_text)[:60]}...\"")
                print(f"{prefix}--------------------------------------------------")


def analyze_pptx_layout(file_path):
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到文件：{file_path}")
        return

    prs = Presentation(file_path)

    # 获取当前 PPT 的版面尺寸大小
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 转换成英寸供直观查看（通常宽屏 PPT 为 13.33" x 7.5"）
    page_width_in = slide_width / Inches(1)
    page_height_in = slide_height / Inches(1)

    print(f"==================================================")
    print(f"📊 PPT 版面分析")
    print(f"   文件名: {os.path.basename(file_path)}")
    print(f"   版面物理大小: 宽 {page_width_in:.2f} 英寸, 高 {page_height_in:.2f} 英寸")
    print(f"==================================================")

    for slide_index, slide in enumerate(prs.slides):
        print(f"\n📂 [第 {slide_index + 1} 页 Slide]")
        process_shapes_relative_position(slide.shapes, slide_width, slide_height, slide_index, prefix="  ")


if __name__ == "__main__":
    # 你的 PPTX 文件路径
    input_ppt = r"C:\Users\H\Desktop\word解析和还原\parser_style\pptx\文件\P19-36_含不可编辑_新材料产品介绍（中越双语)_bilingual (1).pptx"

    analyze_pptx_layout(input_ppt)