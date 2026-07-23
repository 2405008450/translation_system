import re
import math
import statistics
from pptx import Presentation
from pptx.util import Inches, Pt

from pptx.oxml.ns import qn
from common import (estimate_required_height, collect_font_sizes,
                    get_autofit_scale, get_line_spacing_info)

MIN_FONT_SCALE = 0.3        # フォント縮小の下限（不可読/0 を防ぐ）
MAX_FONT_SCALE = 2.0        # フォント拡大の上限（過大化を防ぐ）
DEFAULT_INHERIT_PT = 18.0   # 母版継承で字号が取れない run のフォールバック基準(pt)
MIN_RATIO_TO_SHRINK = 1.15  # model_scale: これ以下の overflow_ratio は誤差とみなし縮小しない
BOTH_MAX_DRIFT_RATIO = 0.3  # both: |new_left-old_left|/old_w や top方向が此の比率を超えたら平移量過大と判定
UNDERFLOW_FILL_MARGIN = 0.10  # underflow拡大時に框内に残す余白率（10%）
MIN_ENLARGE_SCALE = 1.05    # これ未満の拡大は誤差とみなし放大しない（過敏な微調整を防ぐ）

def iter_all_shapes(shape_container):
    for shape in shape_container.shapes:
        yield shape
        if hasattr(shape, 'shapes'):   # GroupShape 和某些容器有 'shapes' 属性
            yield from iter_all_shapes(shape)

def _resolve_base_pt(shape):
    """
    継承字号(run.font.size=None)の run に適用する『基準字号』を推定する。
    框内で明示指定された字号があればその中央値、無ければ既定値を返す。
    """
    sizes = collect_font_sizes(shape)   # 明示指定された run/段落の pt を収集
    return statistics.median(sizes) if sizes else DEFAULT_INHERIT_PT


def clean_text(text):
    """清理字符串中的空白字符、换行符，以便进行模糊匹配"""
    if not text:
        return ""
    return re.sub(r"\s+", "", text)


def print_original_style(slide_index, shape):
    """输出文本框原有的样式和信息"""
    print(f"\n======== [第 {slide_index + 1} 页] 找到目标文本框 ========")
    # 转换为英寸查看
    orig_left = shape.left / Inches(1)
    orig_top = shape.top / Inches(1)
    orig_width = shape.width / Inches(1)
    orig_height = shape.height / Inches(1)

    print(
        f"【原始位置与大小】: Left: {orig_left:.2f}寸, Top: {orig_top:.2f}寸, Width: {orig_width:.2f}寸, Height: {orig_height:.2f}寸"
    )
    print(f"【文本框段落总数】: {len(shape.text_frame.paragraphs)} 个段落")

    # 遍历打印前几个段落的样式作为参考
    print("【原文本框部分内容及样式预览】:")
    for i, para in enumerate(shape.text_frame.paragraphs):
        # 限制输出，防止控制台刷屏
        if i >= 3 and len(shape.text_frame.paragraphs) > 4:
            print(f"  ... 还有 {len(shape.text_frame.paragraphs) - i} 个段落省略未打印")
            break

        para_text = para.text.strip()
        if not para_text:
            continue

        # 获取段落里第一个run的样式作为代表
        font_name = "默认"
        font_size = "默认"
        if para.runs:
            run = para.runs[0]
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = f"{run.font.size.pt} Pt"

        print(
            f"  - 段落 [{i+1}]: \"{para_text[:30]}...\" (字体: {font_name}, 字号: {font_size})"
        )
    print("==================================================\n")


def scale_text_box_fonts(shape, scale_factor):
    """
    遍历文本框中的所有段落和字符块，将其字号按比例缩放。
    ★关键修复：对继承母版字号(run.font.size=None)的 run，先用框内基准字号回退，
      再乘以 scale 并【显式写回】，否则这些文字根本不会被缩放（此前只有显式字号的框能缩）。
    ★同时采集缩放前/后的实际字号(pt)并打印，便于确认字号确实被调整成功。
    """
    base_pt = _resolve_base_pt(shape)   # 继承字号的回退基准
    before_pts, after_pts = [], []      # 改前/改后字号采样(pt)
    for paragraph in shape.text_frame.paragraphs:
        # 1. 段落自带字号（若有）按比例缩放
        if paragraph.font.size:
            before_pts.append(round(paragraph.font.size.pt, 1))
            new_p = round(max(paragraph.font.size.pt * scale_factor, 1.0), 1)
            paragraph.font.size = Pt(new_p)
            after_pts.append(new_p)

        # 2. 逐个 Run 缩放：有显式字号用自身，无则用 base_pt 回退，统一显式写回
        for run in paragraph.runs:
            cur_pt = run.font.size.pt if run.font.size else base_pt
            before_pts.append(round(cur_pt, 1))
            new_size = round(max(cur_pt * scale_factor, 1.0), 1)   # 下限保护(>=1pt)
            run.font.size = Pt(new_size)
            after_pts.append(new_size)

    if before_pts:
        b_lo, b_hi = min(before_pts), max(before_pts)
        a_lo, a_hi = min(after_pts), max(after_pts)
        b_rng = f"{b_lo}pt" if b_lo == b_hi else f"{b_lo}~{b_hi}pt"
        a_rng = f"{a_lo}pt" if a_lo == a_hi else f"{a_lo}~{a_hi}pt"
        print(f"      · 字号 {b_rng} → {a_rng} (scale={scale_factor:.3f}, 共{len(before_pts)}处)")
    else:
        print(f"      · ⚠ 未采集到任何字号(该框可能无文字run)，scale={scale_factor:.3f}")


def _font_scale_by_area(box_w, box_h, req_w, req_h):
    """
    ★フォント縮放率の再定義（面積ベース）:
      文字を s 倍に縮めると 1行の字数∝1/s・行数∝s・行高∝s となり、
      文字が占める『面積』は s^2 に比例する。
      よって、字号不変で必要な面積 A_req=(req_w×req_h) を、
      収めたい枠の面積 A_box=(box_w×box_h) に収めるには:
          A_req × s^2 = A_box  →  s = sqrt(A_box / A_req)
      幅を変えない特殊ケースでは sqrt(box_h/req_h) に一致し、旧・高さ版の一般化になる。
    """
    a_box = max(box_w * box_h, 1e-6)
    a_req = max(req_w * req_h, 1e-6)
    if a_req <= a_box:
        return 1.0
    return max(math.sqrt(a_box / a_req), MIN_FONT_SCALE)


def _font_scale_to_fill(box_w, box_h, req_w, req_h):
    """
    ★underflow(文字が框内で小さすぎ・余白過多)用の『拡大率』:
      文字を s 倍にすると占有面積は s^2 倍。框面積の (1-余白率) までを埋めるよう、
        A_req × s^2 = A_box × (1 - margin)  →  s = sqrt(A_box(1-margin) / A_req)
      上限 MAX_FONT_SCALE で頭打ち。框に余裕が無ければ 1.0(据え置き)。
    """
    a_box = max(box_w * box_h, 1e-6) * (1.0 - UNDERFLOW_FILL_MARGIN)
    a_req = max(req_w * req_h, 1e-6)
    if a_req >= a_box:
        return 1.0
    return min(math.sqrt(a_box / a_req), MAX_FONT_SCALE)


def fit_textbox(shape, target, mode="shrink", max_expand=None, xform=(0.0, 0.0, 1.0, 1.0),
                allow_enlarge=True):
    """
    パイプライン向け: 溢れた文本框を収める。
    target = {"left","top","width","height"}（英寸, スライド絶対座標）はモデルが返す
             『理想框』（＝現フォントのまま全文字を収めるのに必要な位置・大きさ）。

    xform = (ax, ay, sx, sy): この shape の子座標系→スライド絶対座標への affine 変換。
            組合(GroupShape)の外の shape は (0,0,1,1)=恒等でよい。
            ★読み取り(モデル送審)は絶対座標だが、python-pptx への書き戻しは
              子座標系で行う必要があるため、書き込み直前に絶対→子座標へ逆変換する。
              これをしないと組合が二重に変換して左下等へズレる。

    mode:
      "expand" : 枠を target の位置・大きさへ移動＋リサイズし、フォントは変えない
      "shrink" : 枠(元の位置・大きさ)を維持し、面積比 sqrt(A_old/A_req) でフォント縮小
      "both"   : height を max_expand まで拡張(位置/幅は target 採用)、残差を面積比で縮小

    戻り値: フォント縮放率
    """
    E = Inches(1)                      # 914400 EMU / inch
    ax, ay, sx, sy = xform
    sx = sx or 1.0
    sy = sy or 1.0

    # 元の几何を「スライド絶対座標(英寸)」で得る（組合内でも真の見た目位置）
    old_l = (ax + (shape.left or 0) * sx) / E
    old_t = (ay + (shape.top or 0) * sy) / E
    old_w = ((shape.width or 0) * sx) / E
    old_h = ((shape.height or 0) * sy) / E

    def _set_abs_geometry(l, t, w, h):
        """絶対座標(英寸)を、この shape の子座標系(EMU)へ逆変換して書き戻す。"""
        shape.left = int(round((l * E - ax) / sx))
        shape.top = int(round((t * E - ay) / sy))
        shape.width = int(round((w * E) / sx))
        shape.height = int(round((h * E) / sy))

    # ★ underflow（框が十分大きく余白過多・文字が小さすぎ、拡大しても遮挡/重疊なしと
    #   モデルが確認済み）: 溢出モードとは独立に、枠は動かさずフォントだけ拡大する。
    #   overflow と underflow は排他（モデル側で担保）なので、ここで早期 return。
    if target.get("underflow"):
        if not allow_enlarge:
            print("   ↳ [underflow] 检测到富余空间，但放大功能已关闭 → 保持不变")
            return 1.0
        rec = {"text": shape.text_frame.text, "width": old_w,
               "height": old_h, "font_sizes": collect_font_sizes(shape),
               "line_spacing_info": get_line_spacing_info(shape)}
        needed_h = estimate_required_height(rec)
        scale = _font_scale_to_fill(old_w, old_h, old_w, needed_h)
        if scale >= MIN_ENLARGE_SCALE:
            print(f"   ↳ [underflow] 枠維持 W{old_w:.2f} H{old_h:.2f}寸, "
                  f"必要H≈{needed_h:.2f} → フォント放大 {scale:.3f}")
            scale_text_box_fonts(shape, scale)
            return scale
        print(f"   ↳ [underflow] 放大空间不足(scale={scale:.3f}<{MIN_ENLARGE_SCALE}) → 保持不变")
        return 1.0

    if mode == "model_scale":
        # ★ 不改文本框几何。模型给出 overflow_ratio(=文字实际高度/框高度)，
        #   代码按面积平方关系换算字号缩放: scale = sqrt(1 / overflow_ratio)。
        #   模型负责“看超出几倍”，代码负责精确换算（可复算）。
        ratio = target.get("overflow_ratio")
        if not ratio or ratio <= MIN_RATIO_TO_SHRINK:
            # 軽微/誤差レベルの溢れは触らない（正常な框を縮めて崩さないため）
            print(f"   ↳ [model_scale] overflow_ratio={ratio}（≤{MIN_RATIO_TO_SHRINK}）→ 保持不变")
            return 1.0
        scale = max(math.sqrt(1.0 / ratio), MIN_FONT_SCALE)
        print(f"   ↳ [model_scale] 枠不变, overflow_ratio={ratio:.2f} → フォント縮放 {scale:.3f}")
        scale_text_box_fonts(shape, scale)
        return scale

    # 以下模式需要模型给出的理想框几何
    req_l, req_t = target["left"], target["top"]
    req_w, req_h = target["width"], target["height"]
    if req_w <= 0 or req_h <= 0 or old_w <= 0 or old_h <= 0:
        return 1.0

    if mode == "expand":
        _set_abs_geometry(req_l, req_t, req_w, req_h)
        print(f"   ↳ [expand] 枠を移動/リサイズ → L{req_l:.2f} T{req_t:.2f} "
              f"W{req_w:.2f} H{req_h:.2f}寸 (フォント不变)")
        return 1.0

    if mode == "both":
        # # ★ 平移量チェック（安全網）：モデルが座標を大幅に動かしていないか検証。
        # #   |位置差| が元の幅/高さの BOTH_MAX_DRIFT_RATIO を超えたら「大幅平移」と判定し、
        # #   信用せず shrink 相当（元の枠を維持＋面積比フォント縮小）にフォールバックする。
        # drift_l = abs(req_l - old_l) / old_w if old_w > 0 else 0
        # drift_t = abs(req_t - old_t) / old_h if old_h > 0 else 0
        # if drift_l > BOTH_MAX_DRIFT_RATIO or drift_t > BOTH_MAX_DRIFT_RATIO:
        #     print(f"   ⚠ [both] 平移量過大 (dL={drift_l:.2f} dT={drift_t:.2f} > "
        #           f"{BOTH_MAX_DRIFT_RATIO})，判定为不可信几何 → 回退 shrink(枠維持)")
        #     scale = _font_scale_by_area(old_w, old_h, req_w, req_h)
        #     if scale < 0.999:
        #         print(f"   ↳ [both→shrink] 枠維持 W{old_w:.2f} H{old_h:.2f}寸, "
        #               f"必要 W{req_w:.2f} H{req_h:.2f}寸 → フォント縮放 {scale:.3f}")
        #         scale_text_box_fonts(shape, scale)
        #     return scale

        # ★ モデルの4値(new_left/top/width/height)を全部適用（移動＋リサイズ）
        app_h = min(req_h, max_expand) if max_expand else req_h   # 高さは任意で上限cap
        _set_abs_geometry(req_l, req_t, req_w, app_h)

        # 適用後の枠に対し、現フォントで実際に必要な高さを見積もり、
        # まだ収まらなければ面積比 sqrt(A_box/A_req) でフォントを縮小（安全網）。
        rec = {"text": shape.text_frame.text, "width": req_w,
               "height": app_h, "font_sizes": collect_font_sizes(shape),
               "line_spacing_info": get_line_spacing_info(shape)}
        needed_h = estimate_required_height(rec)
        scale = _font_scale_by_area(req_w, app_h, req_w, needed_h)
        if scale < 0.999:
            print(f"   ↳ [both] 4値適用 L{req_l:.2f} T{req_t:.2f} W{req_w:.2f} H{app_h:.2f}寸"
                  f" + フォント縮放 {scale:.3f} (必要H≈{needed_h:.2f})")
            scale_text_box_fonts(shape, scale)
        else:
            print(f"   ↳ [both] 4値適用 L{req_l:.2f} T{req_t:.2f} W{req_w:.2f} H{app_h:.2f}寸 (フォント不变)")
        return scale

    # shrink (デフォルト): 元の枠を維持し、面積比でフォントを縮小して収める
    scale = _font_scale_by_area(old_w, old_h, req_w, req_h)
    if scale < 0.999:
        print(f"   ↳ [shrink] 枠維持 W{old_w:.2f} H{old_h:.2f}寸, "
              f"必要 W{req_w:.2f} H{req_h:.2f}寸 → フォント縮放 {scale:.3f}")
        scale_text_box_fonts(shape, scale)
    return scale





def adjust_textbox_by_text(file_path, target_text, new_left, new_top, new_width, new_height, output_path,
                           font_scale=None):
    prs = Presentation(file_path)
    found = False
    cleaned_target = clean_text(target_text)

    for slide_index, slide in enumerate(prs.slides):
        for shape in iter_all_shapes(slide):
            if shape.has_text_frame:
                cleaned_shape_text = clean_text(shape.text_frame.text)

                if cleaned_target in cleaned_shape_text:
                    found = True
                    print_original_style(slide_index, shape)

                    # 计算宽度的缩放比例（如果用户没指定硬性缩放比例，可以用宽度比来算）
                    if font_scale is None:
                        orig_width = shape.width / Inches(1)
                        # 防止除以 0
                        current_scale = new_width / orig_width if orig_width > 0 else 1.0
                    else:
                        current_scale = font_scale

                    # 3. 重新设置文本框的位置和大小
                    print(
                        f"正在调整文本框至 -> Left: {new_left}寸, Top: {new_top}寸, Width: {new_width}寸, Height: {new_height}寸...")
                    shape.left = Inches(new_left)
                    shape.top = Inches(new_top)
                    shape.width = Inches(new_width)
                    shape.height = Inches(new_height)

                    # 4. 执行字体缩放
                    if current_scale != 1.0:
                        print(f"正在将文本框内的字体按比例 {current_scale:.2f} 进行缩放...")
                        scale_text_box_fonts(shape, current_scale)

    if not found:
        print(f"❌ 未在 PPT 中找到包含指定文本的文本框。")
    else:
        prs.save(output_path)
        print(f"\n✅ 修改完成！新文件已保存至: {output_path}")


# ================= 配置参数并运行 =================
if __name__ == "__main__":
    # 输入你现有的 PPT 文件路径
    input_ppt = r"C:\Users\H\Desktop\word解析和还原\parser_style\pptx\文件\测试布局.pptx"

    # 你想要寻找的指定文本（多行文本）
    # 你想要寻找的指定文本（多行文本）
    target_keyword = """聚合物抗裂砂浆
    """
    # 修改后的新文件名
    output_ppt = "my_presentation_adjusted-2.pptx"

    # 设定新的位置和大小（单位：英寸）
    new_l = 4.62
    new_t = 7.60
    new_w = 2.55
    new_h = 1.16

    # 缩放系数。例如：0.8 表示字号变为原来的 80%。
    # 如果设为 None，程序会自动根据【新宽度 / 旧宽度】的比例自动计算缩放系数。
    fixed_font_scale = 0.690

    # 执行调整
    adjust_textbox_by_text(
        input_ppt, target_keyword, new_l, new_t, new_w, new_h, output_ppt,font_scale=fixed_font_scale
    )