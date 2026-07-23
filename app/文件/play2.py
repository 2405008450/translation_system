import os
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def add_tag_label(slide, tag_text, abs_left_emu, abs_top_emu):
    """
    注釈版pptxに『可視の番号タグ』を焼き込む。
    モデルはこのタグ番号で各文本框を識別 → 構造化JSONでタグ番号を返す →
    pipeline が タグ→uid で確実に対象へ書き戻す。
    位置はスライド絶対座標(EMU)。左上角に赤字の小ラベルを置く。
    """
    w, h = Inches(0.55), Inches(0.28)
    tb = slide.shapes.add_textbox(int(abs_left_emu), int(abs_top_emu), w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_top = 0
    tf.text = tag_text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 0, 0)   # 赤字で視認性UP
    # 塗り＝黄でさらに目立たせる（読み取りやすさのため）
    try:
        tb.fill.solid()
        tb.fill.fore_color.rgb = RGBColor(255, 255, 0)
    except Exception:
        pass
    return tb


def draw_border_and_force_visible(shape):
    """
    最纯粹的标注与全显函数：
    1. 绝不修改任何文本内容、字号、字体、颜色、对齐方式、换行规则及内边距。
    2. 仅通过底层 XML 开启文字的‘垂直与水平溢出渲染’（防止超出边界被截断隐藏）。
    3. 勾勒蓝色边框进行可视化标注。
    """
    tf = shape.text_frame

    # 1. 🔥【底层 XML 强力注入】保持所有样式不动的核心：
    # 无论 PPT 内部排版如何，直接在文件底层命令渲染引擎：“超出的文字部分不准隐藏，强行渲染出来！”
    try:
        bodyPr = tf._txBody.bodyPr
        bodyPr.set('vertOverflow', 'overflow')
        bodyPr.set('horzOverflow', 'overflow')
    except Exception:
        pass

    # 2. 勾勒蓝色可见边框（100%保留原本的物理宽高与位置）
    try:
        line = shape.line
        line.color.rgb = RGBColor(0, 102, 204)  # 蓝色
        line.width = Pt(1.5)                    # 1.5磅粗细
    except Exception:
        pass


def process_shapes_recursively(shapes_collection):
    """
    递归函数：穿透群组(GroupShape)，不放过任何组合里的隐藏文本框
    """
    count = 0
    for shape in shapes_collection:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # 穿透组合
            count += process_shapes_recursively(shape.shapes)
        else:
            # 只要里面包含文字（不论多少），就执行标注和全显处理
            if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
                draw_border_and_force_visible(shape)
                count += 1
    return count


def fix_all_text_boxes_perfectly(file_path, output_path):
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到输入文件：{file_path}")
        return

    prs = Presentation(file_path)
    total_processed = 0

    print(f"开始进行【仅标框与文字全显】的纯净处理...")
    print("=" * 60)

    for slide_index, slide in enumerate(prs.slides):
        processed_in_slide = process_shapes_recursively(slide.shapes)
        total_processed += processed_in_slide
        print(f"📂 [第 {slide_index + 1} 页 Slide] 完成处理，共标记了 {processed_in_slide} 个文本框")

    print("=" * 60)

    if total_processed > 0:
        prs.save(output_path)
        print(f"✅ 处理完毕！未修改任何原文样式，超出的隐藏文字已全显，并已成功标注蓝色边框。")
        print(f"💾 新文件已保存至: {output_path}")
    else:
        print("⚠ 未在 PPT 中找到任何有文字的文本框。")


if __name__ == "__main__":
    # 输入与输出路径
    input_ppt = r"C:\Users\H\Desktop\word解析和还原\parser_style\pptx\文件\pipeline_final.pptx"
    output_ppt = "my_presentation_all_visible_pure.pptx"

    fix_all_text_boxes_perfectly(input_ppt, output_ppt)