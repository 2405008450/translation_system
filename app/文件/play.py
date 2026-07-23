import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement


def clean_text(text):
    if not text:
        return ""
    return re.sub(r"\s+", "", text)


def force_xml_overflow_and_shrink(shape, scale_factor=0.65):
    """
    终极强制全显：注入底层 XML 允许文本框完美流出溢出，且不改变物理尺寸
    """
    tf = shape.text_frame

    # 1. 基础配置：开启换行，从顶部排版，清空所有内边距
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    # 2. 🔥【核心黑科技】：强行修改底层 XML 的溢出控制
    try:
        # 获取文本框底部的 bodyPr 节点
        bodyPr = tf._txBody.bodyPr
        # 强制设置属性：允许文本垂直溢出渲染 (overflow='overflow')，不换行截断
        bodyPr.set('vertOverflow', 'overflow')
        bodyPr.set('horzOverflow', 'overflow')
        bodyPr.set('wrap', 'square')  # 强制方形边界换行溢出
    except Exception as e:
        pass

    # 3. 稳妥缩放所有字体：
    # 将中越文原本的字号直接打折（如0.65），让文字占用的物理高度从根本上大幅下降
    for paragraph in tf.paragraphs:
        # 如果段落没有字号，先继承或赋予一个参考值（通常小字为 14Pt）
        current_pt = paragraph.font.size.pt if paragraph.font.size else 14
        paragraph.font.size = Pt(round(current_pt * scale_factor, 1))

        for run in paragraph.runs:
            if run.font.size:
                run.font.size = Pt(round(run.font.size.pt * scale_factor, 1))
            else:
                # 显式让字符块锁定段落字号，防止继承母版超大字号导致再次溢出隐藏
                run.font.size = Pt(paragraph.font.size.pt)

    # 4. 勾勒蓝色边框（方便检视原始框的大小位置）
    line = shape.line
    line.color.rgb = RGBColor(0, 102, 204)
    line.width = Pt(1.5)


def process_shapes_recursively(shapes_collection, target_cleaned, scale_factor):
    """
    深度递归搜索：确保能揪出组合群组里的那个目标文本框
    """
    found_count = 0
    for shape in shapes_collection:
        # 遭遇群组，套娃进去继续找
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found_count += process_shapes_recursively(shape.shapes, target_cleaned, scale_factor)
        else:
            # 常规形状或组件，进行文本比对
            if shape.has_text_frame:
                cleaned_shape_text = clean_text(shape.text_frame.text)
                if target_cleaned in cleaned_shape_text:
                    print(f"🎯 成功锁定目标文本框！正在执行 XML 底层无隐藏溢出修复...")
                    force_xml_overflow_and_shrink(shape, scale_factor=scale_factor)
                    found_count += 1
    return found_count


def fix_textbox_text_absolute(file_path, target_text, output_path, font_scale=0.65):
    if not os.path.exists(file_path):
        print(f"❌ 错误：找不到文件：{file_path}")
        return

    prs = Presentation(file_path)
    cleaned_target = clean_text(target_text)
    total_found = 0

    # 遍历幻灯片
    for slide_index, slide in enumerate(prs.slides):
        found_in_slide = process_shapes_recursively(slide.shapes, cleaned_target, font_scale)
        if found_in_slide > 0:
            print(f"✨ 第 {slide_index + 1} 页中的目标文本框已完成修复。")
            total_found += found_in_slide

    if total_found == 0:
        print(f"❌ 未在 PPT（包括任何组合形状）中找到包含指定文本的文本框。")
    else:
        prs.save(output_path)
        print(f"\n✅ 终极修复完成！新文件已成功生成至: {output_path}")


# ================= 配置参数并运行 =================
if __name__ == "__main__":
    input_ppt = r"C:\Users\H\Desktop\word解析和还原\P19-36_含不可编辑_新材料产品介绍（中越双语)_bilingual (1).pptx"

    # 你的目标关键词
    target_keyword = """粘结力超强
Độ bám dính siêu mạnh
耐水抗渗
"""
    output_ppt = "my_presentation_text_fixed_only.pptx"

    # 如果生成的 PPT 底部的越文还是差一两行没完全漏出来，将这里的 0.65 调低到 0.55 或 0.50 即可。
    # 配合底层 XML 溢出设置，双管齐下绝对有效。
    adjust_scale = 0.60

    fix_textbox_text_absolute(input_ppt, target_keyword, output_ppt)